# -*- coding: utf-8 -*-
# ---
# jupyter:
#   jupytext:
#     formats: ipynb,py:light
#     text_representation:
#       extension: .py
#       format_name: light
#       format_version: '1.5'
#       jupytext_version: 1.15.0
#   kernelspec:
#     display_name: Python 3 (ipykernel)
#     language: python
#     name: python3
# ---

# # SPMpy 
# * Authors : Dr. Jewook Park at CNMS, ORNL
#     * Center for Nanophase Materials Sciences (CNMS), Oak Ridge National Laboratory (ORNL)
#     * email :  parkj1@ornl.gov
#         
# > **SPMpy** is a python package to analysis scanning probe microscopy (SPM) data analysis, such as scanning tunneling microscopy and spectroscopy (STM/S) data and atomic force microscopy (AFM) images, which are inherently multidimensional. SPMpy exploits recent image processing(a.k.a. Computer Vision) techniques, and utilzes [building blocks](https://scipy-lectures.org/intro/intro.html#the-scientific-python-ecosystem) and excellent visualization tools available in the [scientific python ecosystem](https://holoviz.org/index.html). Many parts are inspired by well-known SPM data analysis programs, for example, [Wsxm](http://www.wsxm.eu/) and [Gwyddion](http://gwyddion.net/). SPMpy is trying to apply lessons from [Fundamentals in Data Visualization](https://clauswilke.com/dataviz/).
#
# >  **SPMpy** is an open-source project. (Github: https://github.com/Jewook-Park/SPMPY )
# > * Contributions, comments, ideas, and error reports are always welcome. Please use the Github page or email parkj1@ornl.gov. Comments & remarks should be in Korean or English. 
#
# * To use SPMpy functions, SPM data() need to be converted as PANDAS DataFrame or Xarray DataSet. 
#
# > * check **SPMpy_fileloading_functions** first.
#

# + [markdown] jp-MarkdownHeadingCollapsed=true
#
#
# -

# # 0.  Import modules

# +

import os
import glob
import numpy as np
import pandas as pd
import scipy as sp
from warnings import warn
from scipy import signal

import math
import skimage
import matplotlib.pyplot as plt
import matplotlib.patches as patches

import seaborn as sns

#from SPMpy_2D_data_analysis_funcs import files_in_folder,img2xr,grid2xr,gridline2xr,gwy_img2df,gwy_df_ch2xr

# some packages may be yet to be installed
try:
     from pptx import Presentation
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named Presentation')
    # !pip install python-pptx  
    from pptx import Presentation
    from pptx.util import Inches, Pt

try:
    import nanonispy as nap
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named nanonispy')
    # !pip install nanonispy
    import nanonispy as nap

try:
    import seaborn_image as isns
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named seaborn-image')
    # #!pip install --upgrade scikit-image == 0.19.0.dev0
    # !pip install --upgrade seaborn-image    
    import seaborn_image as isns

try:
    import xarray as xr
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named xarray')
    # #!pip install --upgrade scikit-image == 0.19.0.dev0
    # !pip install xarray 
    import xarray as xr
    
try:
    import xrft
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named xrft')
    # !pip install xrft 
    import xrft

try:
    import seaborn_image as isns
except ModuleNotFoundError:
    warn('ModuleNotFoundError: No module named seaborn-image')
    # #!pip install --upgrade scikit-image == 0.19.0.dev0
    # !pip install --upgrade seaborn-image    
    import seaborn_image as isns


# -

# ## 1 : plane fit functions
# * plane_fit_y
# * plane_fit_x
# * plane_fit_y_df
# * plane_fit_x_df
# * plane_fit_y_xr
# * plane_fit_x_xr
# * plane_fit_surface
# * plane_fit_surface_df
# * plane_fit_surface_xr

# +
# JW's Functions 
################################################
    # Step 1-2: Import Jewook's new functions 
    #* 1-2-1 : plane fit functions (for np, df, xr)
    #* 1-2-2 : twoD FFT functions (for np, df, xr)
################################################

##################################      
#   Jewook-defined Functions     # 
##################################
#################
#   plane_y_fit
#~~~~~~~~~~~~~~#
# for nd array #
#~~~~~~~~~~~~~~~~~~~~~~~~~#
# numpy  #  nd-in & nd-out#
#~~~~~~~~~~~~~~~~~~~~~~~~~~#
def plane_fit_y(image2D):
    '''
    Parameters
    ----------
    image2D : np.array
        2D scan data  Y (slow scan direction) direction avg move
        
    Returns
    -------
    image_2D_y_fit : np.array
        y plane fit  np results 
    '''
    image_2D_y_fit = (image2D.T- image2D.mean(axis=1)).T
    # transpose image & substract mean (axis =1)
    return image_2D_y_fit
# x fit 
def plane_fit_x(image2D):
    '''
    Parameters
    ----------
    image2D :  np.array
        2D scan data  X (fast scan direction) direction plane fit
    Returns
    -------
    image2D_x_fit : np.array
        x plane fit  np results 
    '''
    image2D_x_fit = image2D -  image2D.mean(axis=0)
    # substract Y direciton mean to X axis 
    return image2D_x_fit
##########################
#~~~~~~~~~~~~~~~#
# for dataframe #
#~~~~~~~~~~~~~~~~~~~~~~~~~~#
# PANDAS #   df-in & df-out#
#~~~~~~~~~~~~~~~~~~~~~~~~~~#
# y fit
def plane_fit_y_df(image2D_df):
    '''
    Parameters
    ----------
    image2D_df : pd.DataFrame input 
        2D scan data  Y (slow scan direction) direction avg move
    Returns
    -------
    image_2D_y_fit_df :pd.DataFrame output
        y plane fit.
    '''
    image2D = image2D_df.to_numpy() # to numpy 
    image_2D_y_fit = plane_fit_y(image2D)
    # use the plane_fit_y function 
    # Transpose the image & substract mean (axis=1) 
    image_2D_y_fit_df = pd.DataFrame(
        image_2D_y_fit,
        index = image2D_df.index,
        columns = image2D_df.columns)
    return image_2D_y_fit_df
#~~~~~~~~~~~~~~#
# x fit 
def plane_fit_x_df(image2D_df):
    '''
    Parameters
    ----------
    image2D_df : pd.DataFrame input 
                2D scan data  X (fast scan direction) direction plane fit
                
    Returns
    -------
    image2D_x_fit_df : pd.DataFrame output
        x plane fit.
    '''
    image2D = image2D_df.to_numpy() # to numpy
    
    image2D_x_fit = plane_fit_x(image2D) 
    # use the plane_fitY function 
    # substract Y direciton mean to X axis 
    image2D_x_fit_df = pd.DataFrame(
        image2D_x_fit,
        index = image2D_df.index, 
        columns = image2D_df.columns)
    # covert  as dataframe
    return image2D_x_fit_df
###################################################
#~~~~~~~~~~~~~~~#
# for xarray    #
#~~~~~~~~~~~~~~~~~~~~~~~~~#
# XARRAY #  xr-in & xr-out#
#~~~~~~~~~~~~~~~~~~~~~~~~~#
def plane_fit_y_xr(xrdata):
    '''
    Parameters
    ----------
    xrdata : xarray  input 
                
    Returns
    -------
    xrdata_prcssd : xarray output
    '''
    xrdata_prcssd = xrdata.copy()
    
    for ch_name in xrdata:
        xrdata_prcssd[ch_name].values = plane_fit_y(xrdata[ch_name])
    return xrdata_prcssd

def plane_fit_x_xr(xrdata):
    '''
    Parameters
    ----------
    xrdata : xarray input 
                
    Returns
    -------
    xrdata_prcssd : xarray output
    '''
    xrdata_prcssd = xrdata.copy()
    for ch_name in xrdata:
        xrdata_prcssd[ch_name].values = plane_fit_x(xrdata[ch_name])
        
    return xrdata_prcssd



########################################



##################################      
#    plane fitted topography     # 
#~~~~~~~~~~~~~~#
# for nd array #
#~~~~~~~~~~~~~~~~~~~~~~~~~#
# numpy  #  nd-in & nd-out#
#~~~~~~~~~~~~~~~~~~~~~~~~~~#
def plane_fit_surface (image2D, order = 2):
    """
    Parameters
    ----------
    image2D : np data 2D
        raw data (topography) .
    order : int (1 or 2), optional
        DESCRIPTION. 
        The default is 2
        2nd order polynomial fitting for the background.

    Returns
    -------
    image2D_plane_sub : np data 2D
        DESCRIPTION.
        after substract the estimated plane from input data

    """
    # regular grid covering the domain of the data
    import scipy as sp
    row_y = image2D.shape[0] # row_y
    col_x = image2D.shape[1] # col_x
    
    x = np.arange(0 , row_y)
    y = np.arange(0 , col_x)
    
    X,Y = np.meshgrid(x,y)
    XX = X.flatten()
    YY = Y.flatten()
    
    data_t = np.array([X.ravel(),Y.ravel(), image2D.ravel()])
    data = np.transpose(data_t) # 각각의 point 위치와 topo를 이용한 2차원 array
    
    #order = 2     # 1: linear, 2: quadratic
    if order == 1:
        # best-fit linear plane
        A = np.c_[data[:,0], data[:,1], np.ones(data.shape[0])]
        C,_,_,_ = sp.linalg.lstsq(A, data[:,2])    # coefficients
        
        # evaluate it on grid
        Z = C[0]*X + C[1]*Y + C[2]
        
        # or expressed using matrix/vector product
        #Z = np.dot(np.c_[XX, YY, np.ones(XX.shape)], C).reshape(X.shape)
    
    elif order == 2:
        # best-fit quadratic curve
        A = np.c_[np.ones(data.shape[0]), data[:,:2], 
                  np.prod(data[:,:2], axis=1), data[:,:2]**2]
        C,_,_,_ = sp.linalg.lstsq(A, data[:,2])
        
        # evaluate it on a grid
        Z = np.dot(np.c_[np.ones(XX.shape), XX, YY, XX*YY,
                         XX**2, YY**2], C).reshape(X.shape)
        image2D_plane_sub = image2D - Z  
        return image2D_plane_sub

########################################################################


def plane_fit_surface_df (image2D_df, order = 2):
    """
    Parameters
    ----------
    image2D_df : PANDAS DataFrame 2D
        raw data (topography) .
    order : int (1 or 2), optional
        DESCRIPTION. 
        The default is 2
        2nd order polynomial fitting for the background.

    Returns
    -------
    image2D_plane_sub_df : PANDAS DataFrame 2D
        DESCRIPTION.
        after substract the estimated plane from input data

    """
    image2D = image2D_df.to_numpy()
    
    image2D_plane_sub =  plane_fit_surface(image2D)        
        
    image2D_plane_sub_df = pd.DataFrame(
    image2D_plane_sub,
    index = image2D_df.index, 
    columns = image2D_df.columns)
    # covert  as dataframe

    return image2D_plane_sub_df

def plane_fit_surface_xr (xrdata, order = 2):
    """

    Parameters
    ----------
    xrdata : Xarray data wtih ch_names
        raw data (topography) .
    order : int (1 or 2), optional
        DESCRIPTION. 
        The default is 2
        2nd order polynomial fitting for the background.

    Returns
    -------
    image2D_plane_sub_df : np data 2D
        DESCRIPTION.
        after substract the estimated plane from input data

    """
    import scipy as sp
    xrdata_prcssd = xrdata.copy()
    
    for ch_name in xrdata:
        image2D = xrdata_prcssd[ch_name].values

        # use the X,Y coords in xrdata 
        X,Y = np.meshgrid(xrdata.X,xrdata.Y)
        XX = X.flatten()
        YY = Y.flatten()

        data_t = np.array([X.ravel(),Y.ravel(), image2D.ravel()])
        data = np.transpose(data_t) # 각각의 point 위치와 topo를 이용한 2차원 array

        #order = 2     # 1: linear, 2: quadratic
        if order == 1:
            # best-fit linear plane
            A = np.c_[data[:,0], data[:,1], np.ones(data.shape[0])]
            C,_,_,_ = sp.linalg.lstsq(A, data[:,2])    # coefficients

            # evaluate it on grid
            Z = C[0]*X + C[1]*Y + C[2]

            # or expressed using matrix/vector product
            #Z = np.dot(np.c_[XX, YY, np.ones(XX.shape)], C).reshape(X.shape)

        elif order == 2:
            # best-fit quadratic curve
            A = np.c_[np.ones(data.shape[0]), data[:,:2], 
                      np.prod(data[:,:2], axis=1), data[:,:2]**2]
            C,_,_,_ = sp.linalg.lstsq(A, data[:,2])

            # evaluate it on a grid
            Z = np.dot(np.c_[np.ones(XX.shape), XX, YY, XX*YY,
                             XX**2, YY**2], C).reshape(X.shape)
            image2D_plane_sub = image2D - Z  

        xrdata_prcssd[ch_name].values = image2D_plane_sub
        
        
    # covert  as xarray data
    
    return xrdata_prcssd

########################################################################






########################################
# mask filtered plaen fit for xr #
########################################
########################################
# xarray columns --> list --> 1st elements 
# masking with condition 
# mean value checkup 
# later 
# -

# ## 2 : TwoD FFT functions
# * twoD_FFT
# * twoD_FFT_df
# * twoD_FFT_xr

# +
#########
# 2D FFT # 
# for numpy & df #
# use np.fft function
#~~~~~~~~~~~~~~~~#
# numpy 
def twoD_FFT(image):
    '''    
    Parameters
    ----------
    image : 2D numpy input 
        Calcuate Fast Fourier Transform.
        # for the np and df, 
        # after fft, pixel size = even number 
        # for line profile at center : use mean value ( center -1  & center +1 )
        # FFT results size =  1/ (pixel size) 
        # topo( 10nm, 512px ) ==> # fft image size  =  512/10  px/nm 
        # 1  pixel size =  1/10 (1/nm)
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #size_x  : image size #step_dx : pixel size 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #FFT size 
        #Grid_3D_xr_nm.X.shape[0]/size_x #Grid_3D_xr_nm.Y.shape[0]/size_y
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        # FFT nX  : # Grid_3D_xr_nm.X/step_dx  # int pixel 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #Grid_3D_xr_nm.X/step_dx*(1/size_x)
        #Grid_3D_xr_nm.Y/step_dy*(1/size_y)
        # integer * unit(1/size_x) 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
    Returns
    -------
    image_fft :  2D numpy output
        To avoid negative value, 
        (add 1 (+1) to fft result) & log.
    '''
    fft = np.fft.fft2(image) ### fft only
    fft_shift = np.fft.fftshift(fft)   # image shift
    image_fft=np.log(np.abs(fft_shift)+1) # + to avoid neg value after log 
    return image_fft


##########################
# PANDAS 
def twoD_FFT_df(image_df):
    """
    Parameters
    ----------
    image_df :  pd.DataFrame input 
        Calcuate Fast Fourier Transform of DataFrame format.
        # for the np and df, 
        # after fft, pixel size = even number 
        # for line profile at center : use mean value ( center -1  & center +1 )
        # FFT results size =  1/ (pixel size) 
        # topo( 10nm, 512px ) ==> # fft image size  =  512/10  px/nm 
        # 1  pixel size =  1/10 (1/nm)
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #size_x  : image size #step_dx : pixel size 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #FFT size 
        #Grid_3D_xr_nm.X.shape[0]/size_x #Grid_3D_xr_nm.Y.shape[0]/size_y
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        # FFT nX  : # Grid_3D_xr_nm.X/step_dx  # int pixel 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
        #Grid_3D_xr_nm.X/step_dx*(1/size_x)
        #Grid_3D_xr_nm.Y/step_dy*(1/size_y)
        # integer * unit(1/size_x) 
        #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
    Returns
    -------
    image_fft_df : pd.DataFrame output 
        (add 1 (+1) to fft result) & log.

    """
    image = image_df.to_numpy()
    image_fft = twoD_FFT(image)
    image_fft_df = pd.DataFrame(image_fft)
    
    return image_fft_df
##########################
#####################################

#####################
# xr 을 위한 fft 함수 지정 
#####################
# 입력은 xr data 
# 출력은 fft 한 xr data 


def twoD_FFT_xr (xr_data,  
                 true_phase=True, 
                 true_amplitude=False, 
                 plus_one = 1,
                 complex_output = False):
    """
    Calcuate Fast Fourier Transform. (with xrft package )
    # in case of xarray data format 
    # XY, coord ==> freq_X, freq_Y
    # xr data (xarray) -> fft -> 
    # auto-calibration  => freq_X & freq_Y 
    
    Parameters
    ----------
    xr_data : Xarray DataSet TYPE
        DESCRIPTION. input data 
    true_phase : Boolean TYPE, optional
        DESCRIPTION. The default is True.
                        xrft package requires it. for ifft 
    true_amplitude : Boolean TYPE, optional
        DESCRIPTION. The default is False.
            for twoD_FFT img show : False -> log-scale output
            for filter& ifft : True -> linear-scale output
                        
    plus_one : INT TYPE, optional
        DESCRIPTION. The default is 1.
         To avoid negative value,  (add 1 (+1) to fft result) & log 
               
    complex_output : Boolean TYPE, optional
        DESCRIPTION. The default is False.
            for filter * ifft : True to maintain imaginary part information

    Returns
    -------
    xr_data_after_fft : TYPE
        DESCRIPTION.   xr_data_after_fft :  xr_data output
    """

    import xrft
    xr_data_after_fft = xr_data.copy() # room for saving fft result
    # load channels in xr_data     
    for ch_map in xr_data:
        if complex_output == False: 
            print('np.log & np.abs after fft')
            xr_data_after_fft[ch_map+'_fft'] = np.log(
                np.abs(xrft.xrft.fft(xr_data[ch_map],
                                true_phase = true_phase,
                                true_amplitude = true_amplitude))+plus_one)
        else: 
            print('complex128 after fft')
            xr_data_after_fft[ch_map+'_fft']  = xrft.xrft.fft(xr_data[ch_map])
    # 'true_phase=True,true_amplitude=False'
    # true amplitude : False ==> for log scale 
    # true phase : True 2D --> for proper fft   
    xr_data_after_fft= xr_data_after_fft.drop_dims(["X", "Y"])
    xr_data_after_fft.attrs = xr_data.attrs 
    # remove  " X","Y"  dimension (not neccessary)
    return xr_data_after_fft


# -

# ## 3. miscellaneous (data mapping functions + )
# > data mapping ( closest, nearest, outliers) 
# > trim additional axes in plot (trim axs) 
# ~~reciprocal_transform_2d~~
# * get_closest
# * find_nearest
# * remove_outlier_df
# * trim_axs (for plot  & remove un-used subaxes ) 
#

# +
############################   
### find nearest index of array for given value
############################# 
#`values` should be sorted
def get_closest(array, values):
    """
    Parameters
    ----------
    array : TYPE : array --> np.arrary()
        DESCRIPTION.
    values : TYPE : target value 
        DESCRIPTION.

    Returns
    -------
    TYPE original array 
        DESCRIPTION.

    """
    #make sure array is a numpy array
    array = np.array(array)
    # get insert positions
    idxs = np.searchsorted(array, values, side="left")
    # find indexes where previous index is closer
    prev_idx_is_less = ((idxs == len(array))|(np.fabs(
        values - array[np.maximum(idxs-1, 0)]) < np.fabs(
            values - array[np.minimum(idxs, len(array)-1)])))
    idxs[prev_idx_is_less] -= 1

    return array[idxs]

#################################################################
def find_nearest(array, values):
    """
    Parameters
    ----------
    array : TYPE : array --> np.arrary()
        DESCRIPTION.
    values : TYPE : target value 
        DESCRIPTION.

    Returns
    -------
    TYPE : original array 
        DESCRIPTION.
    indices : TYPE : index 
        DESCRIPTION.

    """
    array = np.asarray(array)
    # the last dim must be 1 to broadcast in (array - values) below.
    values = np.expand_dims(values, axis=-1) 
    indices = np.abs(array - values).argmin(axis=-1)

    return array[indices],indices
#################################################################


################################################################
# Remove Outliers
################################################################
# Outliers according to Standard devivations 
def remove_outlier_df(target_df):
    """
    Parameters
    ----------
    target_df : TYPE : PANDAS DataFrame
        DESCRIPTION.
        
    Returns
    -------
        outlier references dd.quantile 1/4 *1.5. 3/4*1.5
        
    search_outlier_df : TYPE : PANDAS DataFrame
        DESCRIPTION.  removed outliers
    filterout_df : TYPE : PANDAS DataFrame
        DESCRIPTION. filtered DataFrame

    """
    quartile_1 = target_df.quantile(0.25)# 1/4 percentile
    quartile_3 = target_df.quantile(0.75)# 3/4 percentile
    IQR = quartile_3 - quartile_1# centeral 50% range 
    condition = (target_df < (quartile_1 - 1.5 * IQR)) | (
	target_df > (quartile_3 + 1.5 * IQR))
	# outlier range setting (IQR * 1.5) 
    condition = condition.any(axis=1)
    search_outlier_df = target_df[condition]
	# find outliers
    filterout_df = target_df.drop(search_outlier_df.index, axis=0)
	# remove outliers
    return search_outlier_df,filterout_df
    # quantile 에서 1/4 - 3/4 사이 간격의 1.5 바깥에 해당하는 영역들 제외 
#################################################################################

#######################
# trimming axs for  (grid) figure 
#######################
# remove not necessary axs  accroding to the number of data set 
def trim_axs(axs, N):
    """
    Reduce *axs* to *N* Axes. All further Axes are removed from the figure.
    """
    axs = axs.flat
    for ax in axs[N:]:
        ax.remove()
    return axs[:N]

#######################

##############################################################
# If you need more functions, pleas add
##############################################################


# -


# # 4. image filtering functions (based on skimage.filters) 
#
# > **filters are applied to xr dataset**
# > * Overwrite or not?
# * filter_gaussian_xr
# * filter_median_xr
# * filter_sobel_xr
# * filter_unsharp_mask_xr 
# * filter_diffofgaussians_xr
#     (similar to band pass filter) 
# > Additional functions will be added later 
# > * eg) z_LIX_fNb_xr_gaussian = **filter_gaussian_xr**(z_LIX_fNb_xr, sigma = 3)
# > * eg)z_LIX_fNb_xr_median = **filter_median_xr **(z_LIX_fNb_xr)
# > * eg) z_LIX_fNb_xr_sobel = __filter_sobel_xr__(z_LIX_fNb_xr)
# > * eg) z_LIX_fNb_xr_unsharp_mask = __filte_unsharp_mask_xr__(z_LIX_fNb_xr)
# > * eg) z_LIX_fNb_xr_difference_of_gaussians =  __filter_diffofgaussians_xr__
# > (z_LIX_fNb_xr,low_sigma = 1, high_sigma = 20)
# > * eg) band pass filters (butterworth: skimage >v.0.19) , thresholds, ... 

# +
##############################################################
# gaussian filtering 
##############################################################
def filter_gaussian_xr (xrdata,
                        sigma = 1,overwrite = False):
    '''    
    Parameters
    ----------
    xrdata : xr_data, input 
    sigma  = 1, optional default
    overwrite  defalut = False
    # if overwrite = True :  overwrite original data 
    # if overwrite == False : (keep) original + (new) processed data 
    # processed data: channel name: ch_name+'_gaussian'
    
    Returns
    -------
    xrdata_prcssd :  xrdata after gaussian,  output
    '''       
    xrdata_prcssd=xrdata.copy() # assign new xr data 
    if overwrite == False:
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name+'_gaussian'] = xrdata[ch_name]
            xrdata_prcssd[ch_name+'_gaussian'].values = skimage.filters.gaussian(
                xrdata[ch_name],
                sigma = sigma)
    else : 
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name].values = skimage.filters.gaussian(
                xrdata[ch_name],
                sigma = sigma)
    return xrdata_prcssd
###############################

##############################################################
# median filtering 
##############################################################
def filter_median_xr (xrdata,
                     overwrite = False):
    '''    
    Parameters
    ----------
    xrdata : xr_data, input 
    overwrite  optional, defalut = False
    # if overwrite = True :  overwrite original data 
    # if overwrite == False : (keep) original + (new) processed data 
    # processed data: channel name: ch_name+'_median'
    
    Returns
    -------
    xrdata_prcssd :  xrdata after _median,  output
    '''  
    xrdata_prcssd = xrdata.copy()
    if overwrite == False:
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name+'_median'] = xrdata[ch_name]
            xrdata_prcssd[ch_name+'_median'].values = skimage.filters.median(
                xrdata[ch_name])
    else : 
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name].values = skimage.filters.median(
                xrdata[ch_name])
    return xrdata_prcssd
############################

##############################################################
# sobel filtering 
##############################################################
def filter_sobel_xr (xrdata,
                     overwrite = False ):
    '''    
    Parameters
    ----------
    xrdata : xr_data, input 
    overwrite  optional, defalut = False
    # if overwrite = True :  overwrite original data 
    # if overwrite == False : (keep) original + (new) processed data 
    # processed data: channel name: ch_name+'_sobel'
    
    Returns
    -------
    xrdata_prcssd :  xrdata after _sobel,  output
    ''' 
    xrdata_prcssd = xrdata.copy()
    if overwrite == False:
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name+'_sobel'] = xrdata[ch_name]
            xrdata_prcssd[ch_name+'_sobel'].values = skimage.filters.sobel(
                xrdata[ch_name])
    else : 
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name].values = skimage.filters.sobel(
                xrdata[ch_name])
    return xrdata_prcssd
############################

##############################################################
# unsharp_mask filtering 
##############################################################
def filter_unsharp_mask_xr (xrdata,
                            overwrite = False ):
    '''    
    Parameters
    ----------
    xrdata : xr_data, input 
    overwrite  optional, defalut = False
    # if overwrite = True :  overwrite original data 
    # if overwrite == False : (keep) original + (new) processed data 
    # processed data: channel name: ch_name+'_unsharp_mask'
    
    Returns
    -------
    xrdata_prcssd :  xrdata after _unsharp_mask,  output
    ''' 
    xrdata_prcssd = xrdata.copy()
    if overwrite == False:
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name+'_unsharp_mask'] = xrdata[ch_name]
            xrdata_prcssd[ch_name+'_unsharp_mask'].values = skimage.filters.unsharp_mask(
                xrdata[ch_name])
    else : 
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name].values = skimage.filters.unsharp_mask(
                xrdata[ch_name])
    return xrdata_prcssd
############################
# Test 
#z_LIX_fNb_xr_unsharp_mask = xr_unsharp_mask_filter(z_LIX_fNb_xr)
#isns.imshow(z_LIX_fNb_xr_unsharp_mask.z_fwd_df_unsharp_mask, origin = "lower")
###############################


# difference_of_gaussians filtering 
def filter_diffofgaussians_xr (xrdata,
                              low_sigma = 5, 
                              high_sigma = 12,
                              overwrite = False ):
    '''    
    Parameters
    ----------
    xrdata : xr_data, input 
    a feature enhancement algorithm (DoG)
    subtraction of one Gaussian blurred(high_sigma) version of an original image from another, 
    less blurred (low_sigma) version of the original
    
    overwrite  optional, defalut = False
    # if overwrite = True :  overwrite original data 
    # if overwrite == False : (keep) original + (new) processed data 
    # processed data: channel name: ch_name+'_difference_of_gaussians'
    
    Returns
    -------
    xrdata_prcssd :  xrdata after _difference_of_gaussians,  output
    ''' 
    xrdata_prcssd = xrdata.copy()
    if overwrite == False:
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name+'_difference_of_gaussians'] = xrdata[ch_name]
            xrdata_prcssd[ch_name+'_difference_of_gaussians'].values = skimage.filters.difference_of_gaussians(
                xrdata[ch_name],
                low_sigma = low_sigma,
                high_sigma = 12 )
    else : 
        for ch_name in  xrdata:
            xrdata_prcssd[ch_name].values = skimage.filters.difference_of_gaussians(
                xrdata[ch_name],
                low_sigma = low_sigma,
                high_sigma = 12 )
    return xrdata_prcssd
############################
# Test 
###############################
'''
from skimage.filters import butterworth
##############################################################
# butterworth filtering 
##############################################################
def filter_butterworth_filter (xrdata,cutoff_frequency_ratio=0.005, high_pass=True):
    xrdata_prcssd = xrdata.copy()
    #새로운 data가 들어갈곳 지정
    xrdata_df = xrdata.to_dataframe()
    #원본데이터를 dataframe으로  변환하여 np로 skimage 처리
    # 현재는 xr[ch].values 로 array 로 바로 뽑아냄. 
    for ch_name in  xrdata_df:        
        xrdata_df_procss = pd.DataFrame(
            skimage.filters.butterworth(
                xrdata[ch_name].unstack(), 
                cutoff_frequency_ratio = cutoff_frequency_ratio, 
                high_pass = high_pass)).stack()
        xrdata_df_procss.index  = xrdata_df.index
        xrdata_prcssd[ch_name+'_butterworth'] = xrdata_df_procss.unstack()
    return xrdata_prcssd
############################
# Test 
z_LIX_fNb_xr_butterworthk = xr_butterworth_filter(z_LIX_fNb_xr)
isns.imshow(z_LIX_fNb_xr_butterworth.z_fwd_df_butterworth, origin = "lower")
###############################
'''


# -

# # 5. rank filteres functions 
# > To use "skimage.filters.rank ", input data type should be uint8/uint16
# > * $\to$ image_to_grayscale (convert image as (0,255))
# > * data float32  $\to$ (rescale intensity (range (0,1) $\to$ img_to_ubyte )$\rightarrow$unit8
# > * Caution! <U> grayscale overwriting the channel name  </u>
# * filter_convert2grayscale
# * filter_gs_mean_xr   
#     * input image :  Gray-scale format 
#     * extrace mean value with disk size (default size = 10 px) 
#     
# * filter_gs_substract_mean_xr
#     * substract mean value with disk size (default size = 10 px) 
#     * input image :  Gray-scale format 
#   
#     
# >* eg) xrdata =  **filter_convert2grayscale** (xrdata)
# >* eg) xrdata =  **filter_gs_mean_xr** (xrdata_gs)
# >* eg) xrdata =  **filter_gs_substract_mean_xr** (xrdata)
#

# # 6. Image Thresholding functions 
# > *  Threshold selections + boolean 
# > * threshold_flip = True : boolean T/F selection
# * single threshold values
#     * __threshold_mean_xr__
#     * __threshold_otsu_xr__
#     * __threshold_minimum_xr__
#         * Return threshold value(s) to have __only two peaks__ in the histogram
#     * ~~threshold_li_xr ($\to$ it takes long time to compute)~~
#     * threshold_triangle_xr
#     * threshold_yen_xr
#         * threshold to distinguish areas
#         * Return threshold value based on Yen’s method.
#     * threshold_isodata_xr
#         * Return threshold value(s) based on ISODATA method
#     * __threshold_multiotsu_xr__     
#         * multi threshold values  (using minimum_v in histogram, class number default = 3)
# * array threshold values 
#     * threshold_sauvola_xr (window_size_odd = 15, k =0.2)
#     * __threshold_local_xr__ (blocksize_odd = 15) 
#
# * check the reference functions in scikit image
#     * https://scikit-image.org/docs/stable/auto_examples/applications/plot_thresholding.html#sphx-glr-auto-examples-applications-plot-thresholding-py
#     * https://scikit-image.org/docs/stable/api/skimage.filters.html#skimage.filters.try_all_threshold
#

# + id="po4Jm_loJfhb"
#############################################
# xr data set change to gray scale (0~255)
#############################################

def filter_convert2grayscale(xrdata): 
    """
    convert data values in to grayscale
    to use " skimage.rank.filters "

    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
            input data type is float 32
        
        

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.
        out data type is unsigned 8 bit (0-255)

    """
    xrdata_prcssd = xrdata.copy()
    for ch_name in  xrdata:
        xrdata_prcssd[ch_name].values = skimage.img_as_ubyte(
            skimage.exposure.rescale_intensity(
                xrdata[ch_name])*0.5+0.5)
    return xrdata_prcssd
#############################################
# test example
#z_LIX_fNb_xr_gs = xrdata_to_grayscale(z_LIX_fNb_xr)
#isns.imshow(z_LIX_fNb_xr_gs.z_fwd_df,
#    origin = "lower")
############################################
# now we can do the skimage.rank.filters 
###########################################
#
###########################################
#  apply mean rank.mean to xr_gs
############################################

def filter_gs_mean_xr(xrdata_gs,disk_radious=10): 
    """
    
    apply mean rank.mean to xr_gs
    Input data is unsigned 8 bit 


    Parameters
    ----------
    xrdata_gs : Xarray DataSet TYPE
        DESCRIPTION.
        Input data is unsigned 8 bit (0-255)
        
    disk_radious : INT TYPE, optional
        DESCRIPTION. The default is 10.
        area radius to get mean value 

    Returns
    -------
    xrdata_gs_prcssd : Xarray DataSet TYPE
        DESCRIPTION.
        Xarray DataSet after mean()

    """
    # check! grayscale input ! 
    xrdata_gs_prcssd = xrdata_gs.copy()
    for ch_name in  xrdata_gs:
        xrdata_gs_prcssd[ch_name].values =  skimage.filters.rank.mean(
            xrdata_gs[ch_name],skimage.draw.disk(disk_radious))
    return xrdata_gs_prcssd
#############################################


###########################################
#  apply mean rank.mean to xr_gs
############################################
def filter_gs_substract_mean_xr (xrdata_gs,disk_radious=10): 
    """
    substract the mean value to emphasize atomic details
       
    Input data is unsigned 8 bit 

    Parameters
    ----------
    xrdata_gs : Xarray DataSet TYPE
        DESCRIPTION.
        Input data is unsigned 8 bit (0-255)
    disk_radious : TYPE, optional
        DESCRIPTION. The default is 10.
        area radius to get mean value 

    Returns
    -------
    xrdata_gs_sub_mean :Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_gs_sub_mean = xrdata_gs.copy()
    for ch_name in  xrdata_gs:
        xrdata_gs_sub_mean[ch_name].values = \
        xrdata_gs[ch_name] - skimage.filters.rank.mean(xrdata_gs[ch_name],disk(disk_radious))
    return xrdata_gs_sub_mean

# test

# +
###############################
# threshold selections+ 
#~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
# single threshold values
#~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
# threshold_mean
# threshold_li
# threshold_otsu
# threshold_triangle
# threshold_yen
# ~~~~~~~~~~~~~~~~~~~~~~~~#
# multi threshold values 
# ~~~~~~~~~~~~~~~~~~~~~~~~#
# threshold_multiotsu
# ~~~~~~~~~~~~~~~~~~~~~~~~#
# array threshold values 
# ~~~~~~~~~~~~~~~~~~~~~~~~#
# threshold_sauvola (window_size_odd =15, k =0.2)
# threshold_local (blocksize_odd = 15)+
#~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#


def threshold_mean_xr (xrdata, threshold_flip = True): 
    """
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
     
     use the mean to get threshold value
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy()
    from skimage.filters import threshold_mean
    for ch_name in xrdata :
        threshold_mean_v = threshold_mean(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_mean': threshold_mean_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_mean_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd


'''
def threshold_li_xr (xrdata, threshold_flip = True): 
    xrdata_prcssd = xrdata.copy()
    from skimage.filters import threshold_li
    for ch_name in xrdata :
        print(ch_name)
        threshold_li_v = threshold_li(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_li': threshold_li_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_li_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd
# threshold li take too long time to compute.. 
'''

def threshold_otsu_xr (xrdata, threshold_flip = True):
    """
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
    use the Otsu to get threshold value
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy()
    from skimage.filters import threshold_otsu
    for ch_name in xrdata :
        print(ch_name)
        threshold_otsu_v = threshold_otsu(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_otsu': threshold_otsu_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_otsu_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return  xrdata_prcssd
    

def threshold_triangle_xr (xrdata, threshold_flip = True): 
    """
    INPUT DATA IS supposed to be grayscale 
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
    use the triangle to get threshold value
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data :grayscale (0-255) 
        if not --> covert to gray scale 
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy()
    from skimage.filters import threshold_triangle
    ch_name_xr = [ch_name for ch_name in xrdata]
    if xrdata[ch_name_xr[0]].dtype  != 'uint8' : 
        xrdata_gs = xrdata_to_grayscale(xrdata)
    else: 
        xrdata_gs = xrdata.copy()
    
    for ch_name in xrdata_gs :
        print(ch_name)
        threshold_triangle_v = threshold_triangle(xrdata_gs[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_triangle': threshold_triangle_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata_gs[ch_name] > threshold_triangle_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata_gs[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd 


def threshold_yen_xr (xrdata, threshold_flip = True): 
    """
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
    use the yen to get threshold value
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy() 
    from skimage.filters import threshold_yen
    for ch_name in xrdata :
        print(ch_name)
        threshold_yen_v = threshold_yen(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_yen': threshold_yen_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_yen_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd 
    
    
def threshold_minimum_xr (xrdata, threshold_flip = True):
    xrdata_prcssd = xrdata.copy() 
    """
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
    using minimum_v in histogram 
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    from skimage.filters import threshold_minimum
    for ch_name in xrdata :
        print(ch_name)
        threshold_minimum_v = threshold_minimum(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_minimum': threshold_minimum_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_minimum_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd


def threshold_isodata_xr (xrdata, threshold_flip = True): 
    
    """
    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold
    Return threshold value(s) based on ISODATA method.
    https://scikit-image.org/docs/0.18.x/api/skimage.filters.html
    
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy() 
    from skimage.filters import threshold_isodata
    for ch_name in xrdata :
        print(ch_name)
        threshold_isodata_v = threshold_isodata(xrdata[ch_name].values)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_isodata': threshold_isodata_v }
        # save the threshold value to the attributes 
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_isodata_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd 


def threshold_multiotsu_xr (xrdata, multiclasses = 3):
    """
    
    Generate classes (3) threshold values to divide gray levels in image.
    input data: grayscale 
    
    
    Parameters
    ----------
    xrdata :  Xarray DataSet TYPE
        DESCRIPTION.
         Input data
    multiclasses : Boolean TYPE, optional
        DESCRIPTION. The default is 3.
        number of divisino 

    Returns
    -------
    xrdata_prcssd : TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy() 
    ch_name_xr = [ch_name for ch_name in xrdata]
    if xrdata[ch_name_xr[0]].dtype  != 'uint8' : 
        xrdata_gs = filter_convert2grayscale(xrdata)
    else: 
        xrdata_gs = xrdata.copy()
        
    from skimage.filters import threshold_multiotsu
    for ch_name in xrdata_gs :
        print(ch_name)
        threshold_multiotsu_v = threshold_multiotsu(xrdata_gs[ch_name].values,
                                                    classes = multiclasses)
        xrdata.attrs['threshold'] = {
            ch_name+'_threshold_multiotsu': threshold_multiotsu_v }
        xrdata_prcssd[ch_name].values = np.digitize(
            xrdata_gs[ch_name], bins = threshold_multiotsu_v)
        # output is threshold value array  (len = multiclasses-1)
        # use to select (np.digitize(image, bins=thresholds))
    return   xrdata_prcssd  


def threshold_sauvola_xr (xrdata, window_size_odd = 15, threshold_flip = True): 
    """
    Applies Sauvola local threshold to an array. Sauvola is a modification of Niblack technique.

    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold


    Parameters
    ----------
    xrdata :  Xarray DataSet TYPE
        DESCRIPTION.
        Input data
    window_size_odd : Integer TYPE, optional
        DESCRIPTION. The default is 15.
        it should be an odd number 
        rectangular area selection for find local threshold value 
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd : Xarray DataSet TYPE
        DESCRIPTION.

    """
    xrdata_prcssd = xrdata.copy() 
    from skimage.filters import threshold_sauvola
    for ch_name in xrdata :
        print(ch_name)
        threshold_sauvola_v = threshold_sauvola(xrdata[ch_name].values, window_size = window_size_odd)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_sauvola': threshold_sauvola_v }
        # output is threshold value array (image size)  
        # each pixel has separate threshold value
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_sauvola_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return xrdata_prcssd 

def threshold_local_xr (xrdata, block_size_odd = 15, threshold_flip = True): 
    """
    Applies Local threshold to an array.
    Compute a threshold mask image based on local pixel neighborhood.

    apply threshold to distinguish areas 
    if  threshold_flip = True
        1: larger than threshold
        0: smaller than threshold


    Parameters
    ----------
    xrdata :  Xarray DataSet TYPE
        DESCRIPTION.
    block_size_odd : Integer TYPE, optional
        DESCRIPTION. The default is 15.
         it should be an odd number 
        rectangular area selection for find local threshold value 
    threshold_flip : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        to assign area w.r.t threshold value 

    Returns
    -------
    xrdata_prcssd :  Xarray DataSet TYPE
        DESCRIPTION.

    """    
    xrdata_prcssd = xrdata.copy() 
    from skimage.filters import threshold_local
    for ch_name in xrdata :
        print(ch_name)
        threshold_local_v = threshold_local(xrdata[ch_name].values,  block_size =  block_size_odd)
        xrdata.attrs['threshold'] = {ch_name+'_threshold_local': threshold_local_v }
        # output is threshold value array (image size)  
        # each pixel has separate threshold value
        xrdata_prcssd_ch_mask = xrdata[ch_name] > threshold_local_v
        if threshold_flip == True: 
            xrdata_prcssd_ch_mask = xrdata_prcssd_ch_mask
        else : 
            xrdata_prcssd_ch_mask = ~ xrdata_prcssd_ch_mask
        xrdata_prcssd[ch_name] = xrdata[ch_name].where(xrdata_prcssd_ch_mask)
    return  xrdata_prcssd
    

    
# test
# all works well except li.. (it takes too long time) 

# -


# # 7. Image equalization functions 
# *  **rescale_intensity_xr** function
#     * image_rescale : robust in isns plot ( eg, plot image percentile 0.02~ 0.98)
#     * image_rescale to remove outliers
# *  **equalize_hist_xr** function
#     * histogram_equalization  : to increase global contrast, re-distribute intensity -> lost highest or lowest infomation 
#     *  adatptive_equalization : contrast limited adaptive histogram equalization (CLAHE) -> histo_equalization with small blocks. 
# * **adaptive_equalize_hist_xr** function 
#     * adjust global histogram $ \to$  **(caution)** it can distort the Z information 
#     
# > * <u> to confirm the Z distortion, check the line profile </u>
#

# +

def rescale_intensity_xr (xrdata, percentile = (2,98)): 
    """
    

    Parameters
    ----------
    xrdata : xarray data
        input dataset .
    percentile : tuple TYPE, optional
        DESCRIPTION. The default  equalization range is (2,98) [%].

    Returns
    -------
    xrdata_prcssd : xarray data
        processed xarrary data.

    """
    xrdata_prcssd =  xrdata.copy()
    from skimage import exposure
    for ch_name in xrdata :
        print(ch_name)
        p_low, p_high = np.percentile(xrdata[ch_name].values, percentile)
        rescaled_image = exposure.rescale_intensity(xrdata[ch_name].values,
                                                    in_range = (p_low, p_high))
        xrdata_prcssd[ch_name].values = rescaled_image
    return xrdata_prcssd



def  equalize_hist_xr(xrdata, mask = None): 
    """
    

    Parameters
    ----------
    xrdata : xarray data
        input dataset .
    mask : boolean data array  TYPE, optional
        the same shape         
        The default is None.
        follow the skimage equalize_hist function description

    Returns
    -------
    xrdata_prcssd :xarray data
        processed xarrary data.

    """
    # mask :  ndarray of bools or 0s and 1s, optional
    xrdata_prcssd =  xrdata.copy()
    from skimage import exposure
    for ch_name in xrdata :
        print(ch_name)
        eq_hist_image = exposure.equalize_hist(xrdata[ch_name].values,
                                                    mask = mask)
        xrdata_prcssd[ch_name].values = eq_hist_image
    return xrdata_prcssd



def  adaptive_equalize_hist_xr(xrdata, clip_limit = 0.03): 
    """
    

    Parameters
    ----------
    xrdata : xarray data
        input dataset .
    clip_limit : float, optional
        DESCRIPTION. The default is 0.03.
        contrast limited adaptive histogram equalization (CLAHE)
        -> histo_equalization with small blocks.
        follow the skimage equalize_adapthist function description

    Returns
    -------
    xrdata_prcssd :  :xarray data
        processed xarrary data.

    """
    # mask :  ndarray of bools or 0s and 1s, optional
    xrdata_prcssd =  xrdata.copy()
    from skimage import exposure
    for ch_name in xrdata :
        print(ch_name)
        adap_eq_hist_image = exposure.equalize_adapthist(xrdata[ch_name].values,
                                                    clip_limit = clip_limit)
        xrdata_prcssd[ch_name].values = adap_eq_hist_image
    return xrdata_prcssd


# + [markdown] id="BRKASaSWwLa8"
# #  8. Plot 2D images (xr data, seaborn-image)
# > Seaborn-image package 
# > * check (https://seaborn-image.sarthakjariwala.com/en/latest/index.html)
#
# * **plot_2D_xr** function
#     * for real space image plot (z fwd bwd + LDOS fwd/bwd)
# > image structure 
# | real space (**r**) |  real space (**r**)|  
# | ---| ---|    
# |topo (fwd)|topo(bwd)|  
# |LIX (fwd) |LIX (fwd)|  
#
# * **plot_2D_xr_fft** function
#     * for k-space image plot (z fwd bwd-2D FFT + LDOS fwd/bwd-2D FFT)
# > image structure 
# | momentum space (**k**) | momentum space (**k**)|
# | ---| ---|
# | topo_fft (fwd)|topo_fft (bwd)| 
# |  LIX_fft  (fwd)|   LIX_fft  bwd)|

# +
def plot_2D_xr(xrdata,
                 ncols = 2,
                 figsize = (6,5),
                 colormap4z = 'bone',
                 colormap4LDOS = 'bwr',
                 save_figure = False): 
    """
    display xr data array channels 
    'z'  'LIX'
    use isns. 'nm' unit, 

    Parameters
    ----------
    xrdata : xarray dataset TYPE
        DESCRIPTION.
    ncols : integer TYPE, optional
        number of columns to grid image display.
        The default is 2.
    figsize : tuple TYPE, optional
        image size. The default is (6,5).
    colormap4z : matplotlib colormap str TYPE, optional
        colormap for z map. The default is 'bone'.
    colormap4LDOS :  matplotlib colormap str , optional
        colormap for LDOS (LIX) map.. The default is 'bwr'.
    save_figure :  boolean TYPE, optional
        If it is true= use the xrdata title as a file name.
        The default is False.

    Returns
    -------
    fig : image show  TYPE
        need to chekc the figure saving 
        

    """
    ####################################
    isns.set_image(cmap = 'viridis', origin ="lower")
    fig,axes = plt.subplots(ncols = ncols,
                            nrows = len(xrdata.data_vars)//ncols+1,
                            figsize = figsize)
    axs = axes.ravel()
    axs = Trim_axs(axs, len(xrdata.data_vars))
    # triming first with respect to # of data channel 
    isns.set_image(cmap = 'viridis',origin='lower')   #  set image  direction
    for i,i_channels in enumerate(xrdata):
        isns_channels= i_channels+'_isns'
        if 'LIX' in i_channels:
            print(isns_channels)
            colormap4img = colormap4LDOS
            colrbar_label = 'dI/dV [A/V]'
        elif 'z' in i_channels:
            print(isns_channels)
            colormap4img = colormap4z
            colrbar_label = 'z [m]'
        else: 
            colormap4img = 'viridis'
            colrbar_label = i_channels
            # colormap setting 
        isns_channels = isns.imshow(xrdata[i_channels],
                                    cmap = colormap4img,
                                    ax =  axs[i],
                                    cbar = True,
                                    dx = xrdata.X_spacing*1E9,
                                    units = "nm",
                                    robust = True,
                                    cbar_label  = colrbar_label,
                                    fontsize = 'small' )
        axs[i].set_title(i_channels, loc='left', fontsize = 'small')

        #################################
            # 4channel 2x2 
    fig.suptitle(xrdata.title , fontsize = 'small', position=(0.5, 1.0-0.05) )
    fig.tight_layout()
    if save_figure == True:
        fig.savefig(xrdata.title.split('\n')[0]+'.png',bbox_inches = 'tight')
    else: pass
    isns.reset_defaults()
    plt.show()
    return fig
##################################################



#####################################################################
def plot_2D_xr_fft(xrdata_fft,
                     ncols = 2, 
                     figsize =  (6,5), 
                     zoom_in_fft = True,
                     zoom_in_quantile = 1/4,
                     colormap4z = 'Greys',
                     colormap4LDOS = 'Blues'): 
    """
    display xr DATA Array channels after FFT 
    (k-space)
    default cmap = 'viridis'
    # use isns. 'nm' unit, 

    Parameters
    ----------
    xrdata_fft : Xarray DataSet  TYPE
        xrdata_fft  for k space plot.
    ncols : Integer TYPE, optional
        number of columns to grid image display.
        The default is 2.
    figsize : tuple TYPE, optional
        out put figure size. The default is (6,5).
    zoom_in_fft : Boolean TYPE, optional
        check the zoomin or not. The default is True.
    zoom_in_quantile : float TYPE, optional
        zoomin ratio compare to original size.
        The default is 1/4.
    colormap4z : matplotlib colormap str TYPE, optional
        colormap for z map fft. The default is 'Greys'.
    colormap4LDOS : matplotlib colormap str TYPE, optional
        colormap for LDOS map fft. The default is 'Blues'.

    Returns
    -------
    fig : image show  TYPE
        check the saving figure results. 
        

    """
    # aspect ratio check 
    [size_x, size_y] = xrdata_fft.image_size
    step_dx = size_x/len(xrdata_fft.freq_X)
    step_dy = size_y/len(xrdata_fft.freq_Y)
    
    scan_aspect_ratio = (size_x/len(xrdata_fft.freq_X))/(size_y/len(xrdata_fft.freq_Y))
    if scan_aspect_ratio == 1: 
        cbar_show = True
    else:
        cbar_show = False
    #######################################
    isns.set_image(origin ="lower")
    #  set image  direction
    fig,axes = plt.subplots(ncols = ncols,
                            nrows = len(xrdata_fft.data_vars)//ncols+1,
                            figsize = figsize)
    axs = axes.ravel()
    axs = Trim_axs(axs, len(xrdata_fft.data_vars))
    # triming first with respect to # of data channel 
    for i,i_channels in enumerate(xrdata_fft):
        isns_channels= i_channels+'_isns'
        if 'LIX' in i_channels:
            #print(isns_channels)
            colormap4img = colormap4LDOS
            #colrbar_label = 'dI/dV [A/V]'
        elif 'z' in i_channels:
            #print(isns_channels)
            colormap4img = colormap4z
            #colrbar_label = 'z [m]'
        else: 
            colormap4img = 'viridis'
            colrbar_label = i_channels
            pass #colormap4img = 'viridis'
            # colormap setting 
        if i_channels.endswith('fft'): # fft channels 
            dx = (1/size_x)*1E-9
            units="1/nm"
            dimension  = "si-reciprocal"
        else : 
            dx=step_dx*1E9
            units="nm"
            dimension  = "si"
        #print(isns_channels)
        isns_channels = isns.imshow(xrdata_fft[i_channels],
                                    ax =  axs[i],
                                    cbar = cbar_show,
                                    cmap = colormap4img,
                                    dx = dx,
                                    units= units,
                                    dimension = dimension,
                                    robust = True)
        axs[i].set_title(i_channels, loc ='left', fontsize = 'small')
        #################################
        if scan_aspect_ratio != 1: 
            fig.colorbar(isns_channels.get_children()[-2],  
                        fraction = 0.045, 
                        ax = axs[i])  
        ################################
        # colorbar setting separately
        ################################
        
        ##########
        # zoom in FFT with respect to the total fft image size         
        ## for zoom in 
        if zoom_in_fft == True:
            print("zoom_in_fft")
            # center position (0,0) -> (0,0) in (freq_X,freq_Y )-> find index
            # np.abs(z_LIX_fNb_xr_fft.freq_X).argmin()
            # x_lower_limit => pixel number
            x_lower_limit = np.abs(xrdata_fft.freq_X 
                                   - np.quantile(xrdata_fft.freq_X.values,
                                                 0.5  - zoom_in_quantile)).argmin()
            x_upper_limit = np.abs(xrdata_fft.freq_X 
                                   -  np.quantile(xrdata_fft.freq_X.values,
                                                  0.5  + zoom_in_quantile)).argmin()\
            # +- zoom_in_quantile ratio         
            print (x_lower_limit,x_upper_limit)
            axs[i].set_xbound([x_lower_limit,x_upper_limit])
            axs[i].set_ybound([x_lower_limit,x_upper_limit])    
        ## for zoom in )
    fig.suptitle(xrdata_fft.title+'(==> 2D FFT)',
                 fontsize = 'small',
                 position=(0.5, 1.0-0.05) )
    fig.tight_layout()
    isns.reset_defaults()
    plt.show()
    
    return fig

#############################################


# -

# # 9. Line profiles (xr DataSet) 
#
# > ### line profile plot of only 1 channle 
# * **line_profile_xr_GUI**
#     * GUI input (__2__ points) (start & end pt)
#     * input parameter 
#         * ch_N = 0,1,2,3 
#             * (0: z_fwd_df,1: z_bwd_df,2: LIX_fwd_df,3: LIX_bwd_df )
#         * profile_width = (default = 3 , averaged )  
#         * channel 3points input -->  line profile btw 1,2 points. 
#     * return :
#         * l_pf_start,l_pf_end,fig
#         * To save the image, use the third term ( _,_, figure)
#     
# * **line_profile_xr**
#     * the same as __line_profile_xr_GUI__
#     * use the given point infomation without GUI input
#     * return fig only. 
#
#
#
# > ### line profile plot of 2 channles
#
# * **line_profile2_xr_GUI**
#     * GUI input (__2__ points) (start & end pt)
#     * input parameter 
#         * ch_N = 0,1,2,3 
#             * (0: z_fwd_df,1: z_bwd_df,2: LIX_fwd_df,3: LIX_bwd_df )
#         * profile_width = (default = 3 , averaged )  
#         * channel 3points input -->  line profile btw 1,2 points. 
#
#     * return :
#         * l_pf_start,l_pf_end,fig
#         * To save the image, use the third term ( _,_, figure)
#         
# * **line_profile2_xr**     
#     * the same as __line_profile2_xr_GUI__
#     * use the given point infomation without GUI input
#     * return fig only. 

# +

def line_profile_xr_GUI(xrdata, ch_N = 0, profile_width = 3):
    """
    # GUI input

    Parameters
    ----------
    xrdata : xarray  DataSet TYPE
        DESCRIPTION. : input data 
    ch_N : Integer TYPE, optional
            call the first Data Array in the Xarray DataSet
        DESCRIPTION. The default is 0.
            Usually the z_fwd_df(or df_fft) 
    profile_width : Integer TYPE, optional
        DESCRIPTION. The default is 3.
            line width (perpendicular to the line ) for profile value 
            use the mean value of perpendicular 3 points 

    Returns
    -------
    l_pf_start : index of the field of view  TYPE
        DESCRIPTION. : starting point of line profile (cx,ry)
    l_pf_end : index of the field of view  TYPE 
        DESCRIPTION.: ednd point of line profile (cx,ry)
    fig    : fig  object TYPE
        DESCRIPTION. : figure object 

    """
    # %matplotlib qt5
    [size_x, size_y] = xrdata.image_size
    # change the matplotlib backend : inline --> qt5  
    # Open a new window 
    for ch_i, ch_name in enumerate(xrdata):
        if ch_i == ch_N :  #  use the first channel image
            fig,axs = plt.subplots (nrows=1,ncols=1, figsize = (6,6))
            isns.imshow(xrdata[ch_name].values,ax = axs, robust=True, origin ="lower", cmap ='copper')
            l_pf_point_2 = fig.ginput(2) 
            # tuple of 2 points 
            plt.show()
            ch_name_to_show = ch_name
        else :
            pass
    print(l_pf_point_2) 
     
    #######################################################################
    # check! l_pf_point_2 is tuple 
    #####################################################################
    l_pf_point_2_np = np.array(l_pf_point_2)
    #####################################################################
    l_pf_start  = l_pf_point_2_np[0].astype(int) # cx, ry output idx
    l_pf_end  =  l_pf_point_2_np[1].astype(int) # cx, ry output idx

    # %matplotlib inline
    # come back to the inline backend. 
    
    l_pf_length = sp.spatial.distance.euclidean(l_pf_start,l_pf_end)
    l_pf_start_end_pt = [*zip(l_pf_start,l_pf_end)]
    scan_aspect_ratio = (size_x/len(xrdata.X))/(size_y/len(xrdata.Y))
    
    l_pf = skimage.measure.profile_line(xrdata[ch_name_to_show].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
 
    fig,axs = plt.subplots (nrows=2,ncols=1, figsize = (5,8))

    isns.imshow(xrdata[ch_name_to_show].values,
                robust=True,
                origin ="lower",
                ax = axs[0], cmap ='copper')
    axs[0].arrow(l_pf_start[0],
                 l_pf_start[1],
                 l_pf_end[0]-l_pf_start[0],
                 l_pf_end[1]-l_pf_start[1],
                 width = 2,
                 color = 'red')
    axs[0].annotate(xrdata.title,
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    '''
    if scan_aspect_ratio != 1: 
        fig.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) '''
    axs[1].plot(np.linspace(0,
                            l_pf_length, 
                            len(l_pf)
                           )*xrdata.attrs['X_spacing']*1E9,
                l_pf*1E9)
    axs[1].grid(True)
    axs[1].set_xlabel( "length (nm)" , fontsize = 12)
    axs[1].set_ylabel( "height (nm)" , fontsize = 12)
    

    axs[1].annotate('Line Profile',
            xy=(0, 1+0.1), xycoords='axes fraction',
            horizontalalignment='left', verticalalignment='top',
            fontsize='medium')
    plt.tight_layout()
    
    fig.tight_layout()
    isns.reset_defaults()
    plt.show()
    
    #fig.savefig('fig1.pdf')
    # save first, then show. 
    #plt.show()
    # after plt.show(), figure is reset, 
        
    return l_pf_start,l_pf_end,fig
# test 



def line_profile_xr(xrdata, l_pf_start, l_pf_end, ch_N = 0, profile_width = 3):
    """
    # Without  GUI input

    Parameters
    ----------
    xrdata : xarray  DataSet TYPE
        DESCRIPTION. : input data 
    ch_N : Integer TYPE, optional
            call the first Data Array in the Xarray DataSet
        DESCRIPTION. The default is 0.
            Usually the z_fwd_df(or df_fft) 
    profile_width : Integer TYPE, optional
        DESCRIPTION. The default is 3.
            line width (perpendicular to the line ) for profile value 
            use the mean value of perpendicular 3 points 
    l_pf_start: np.array([cx,ry])
            astype(int), shape = (2,)
            # cx, ry  idx
    l_pf_end:  np.array([cx,ry])
            astype(int), shape = (2,)
            # cx, ry  idx
    Returns
    -------
    fig    : fig  object TYPE
        DESCRIPTION. : figure object 

    """
    
    [size_x, size_y] = xrdata.image_size
    for ch_i, ch_name in enumerate(xrdata):
        if ch_i == ch_N :  #  use the first channel image
            ch_name_to_show = ch_name
        else :
            pass
    # %matplotlib inline
    # come back to the inline backend. 
    l_pf_start = l_pf_start
    l_pf_end = l_pf_end
    
    l_pf_length = sp.spatial.distance.euclidean(l_pf_start,l_pf_end)
    l_pf_start_end_pt = [*zip(l_pf_start,l_pf_end)]
    scan_aspect_ratio = (size_x/len(xrdata.X))/(size_y/len(xrdata.Y))
    
    l_pf = skimage.measure.profile_line(xrdata[ch_name_to_show].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
 
    fig,axs = plt.subplots (nrows=2,ncols=1, figsize = (5,8))

    isns.imshow(xrdata[ch_name_to_show].values,
                robust=True,
                origin ="lower",
                ax = axs[0])
    axs[0].arrow(l_pf_start[0],
                 l_pf_start[1],
                 l_pf_end[0]-l_pf_start[0],
                 l_pf_end[1]-l_pf_start[1],
                 width = 2,
                 color = 'copper')
    axs[0].annotate(xrdata.title,
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    
    if scan_aspect_ratio != 1: 
        fig.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) 
    axs[1].plot(np.linspace(0,
                            l_pf_length, 
                            len(l_pf)
                           )*xrdata.attrs['X_spacing']*1E9,
                l_pf*1E9)
    axs[1].grid(True)
    axs[1].set_xlabel( "length (nm)" , fontsize = 12)
    axs[1].set_ylabel( "height (nm)" , fontsize = 12)
    

    axs[1].annotate('Line Profile',
            xy=(0, 1+0.1), xycoords='axes fraction',
            horizontalalignment='left', verticalalignment='top',
            fontsize='medium')
    plt.tight_layout()
    
    fig.tight_layout()
    isns.reset_defaults()
    plt.show()
    
    #fig.savefig('fig1.pdf')
    # save first, then show. 
    #plt.show()
    # after plt.show(), figure is reset, 
        
    return fig
# test 



# +
def line_profile2_xr_GUI(xrdata, ch_N = [0,2], profile_width = 3):
    """
    

    Parameters
    ----------
    xrdata : xarray  DataSet TYPE
        DESCRIPTION. : input data 
    ch_N : Integer list TYPE, optional
            use the channel numbers for line profile
            in the Data Array in the Xarray DataSet
            DESCRIPTION. The default is 0,2.
            the z_fwd_df & LIX_fwd_df
    profile_width : Integer TYPE, optional
            The default is 3.
            line width (perpendicular to the line ) for profile value 
            use the mean value of perpendicular 3 points 

    Returns
    -------
    l_pf_start : index of the field of view  TYPE
        DESCRIPTION. : starting point of line profile (cx,ry)
    l_pf_end : index of the field of view  TYPE 
        DESCRIPTION.: ednd point of line profile (cx,ry)
    fig1   : fig object  TYPE
        DESCRIPTION. : figure (line on figures)
    fig2    : fig object TYPE
        DESCRIPTION. : line profile figure (line profiles ) 

    """
    # %matplotlib qt5
    [size_x, size_y] = xrdata.image_size
    # change the matplotlib backend : inline --> qt5  
    # Open a new window 
    ch_names_for_lp = []
    for ch_i, ch_name in enumerate(xrdata):
        if ch_i in ch_N : 
            ch_names_for_lp.append(ch_name)
            # check channel names first 
            
            #  use the first channel image to for 2 point in line profile
            if ch_i == ch_N[0]:
                fig,axs = plt.subplots (nrows=1,ncols=1, figsize = (6,6))
                isns.imshow(xrdata[ch_name].values,ax = axs, robust=True,origin ="lower")
                l_pf_point_2 = fig.ginput(2) 
                # tuple of three points 
                plt.show()
            else: pass 
        
        else :
            pass
    print(l_pf_point_2) 

    print(ch_names_for_lp) 

    #######################################################################
    # check! l_pf_point_2 is tuple 
    #####################################################################
    l_pf_point_2_np = np.array(l_pf_point_2)
    #####################################################################
    l_pf_start  = l_pf_point_2_np[0].astype(int) # cx, ry output idx
    l_pf_end  =  l_pf_point_2_np[1].astype(int) # cx, ry output idx

    # %matplotlib inline
    # come back to the inline backend. 
    
    l_pf_length = sp.spatial.distance.euclidean(l_pf_start,l_pf_end)
    l_pf_start_end_pt = [*zip(l_pf_start,l_pf_end)]
    scan_aspect_ratio = (size_x/len(xrdata.X))/(size_y/len(xrdata.Y))
    
    l_pf1 = skimage.measure.profile_line(xrdata[ch_names_for_lp[0]].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
    
    l_pf2 = skimage.measure.profile_line(xrdata[ch_names_for_lp[1]].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
 
    fig1,axs = plt.subplots (nrows=1,ncols=2, figsize = (6,3))
    isns.imshow(xrdata[ch_names_for_lp[0]].values, 
                robust=True,
                origin ="lower", 
                cmap = 'bone',
                ax = axs[0])
    axs[0].arrow(l_pf_start[0],
             l_pf_start[1],
             l_pf_end[0]-l_pf_start[0],
             l_pf_end[1]-l_pf_start[1],
             width = 3,
             color = 'tab:red')
    axs[0].annotate(ch_names_for_lp[0],
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    '''
    if scan_aspect_ratio != 1: 
        fig1.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) '''
    # image =  0 :( ==> xrdata[ch_names_for_lp[0]]  )
    isns.imshow(xrdata[ch_names_for_lp[1]].values,
                robust=True,
                origin ="lower",
                cmap = 'bwr',
                ax = axs[1])
    
    axs[1].arrow(l_pf_start[0],
         l_pf_start[1],
         l_pf_end[0]-l_pf_start[0],
         l_pf_end[1]-l_pf_start[1],
         width = 3,
         color = 'tab:green')
    axs[1].annotate(ch_names_for_lp[1],
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    '''
    if scan_aspect_ratio != 1: 
        fig1.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) '''
    fig1.suptitle(xrdata.title,
                  fontsize = 'small',
                  position=(0.5, 1.0+0.001) )
    fig1.tight_layout()
    isns.reset_defaults()
    plt.show()    
    
    
    
    fig2,axs = plt.subplots (figsize = (4,2))
    axs.plot(np.linspace(0,
                            l_pf_length, 
                            len(l_pf1)
                           )*xrdata.attrs['X_spacing']*1E9,
                l_pf1*1E9,
            c = 'tab:red')
    axs.grid(True)
    plt.xticks(fontsize =  'small')
    plt.xticks(fontsize =  'small')
    axs.set_xlabel( "Length (nm)" , fontsize = 'small')
    axs.tick_params(axis = 'y', colors = 'red')
    axs.set_ylabel( "Height (nm)" , fontsize =  'small',c = 'tab:red')
    
    axs1 = axs.twinx()
    axs1.plot(np.linspace(0,
                        l_pf_length, 
                        len(l_pf1)
                       )*xrdata.attrs['X_spacing']*1E9,
              l_pf2*1E9,
             c = 'tab:green')    
    axs1.set_ylabel( "dI/dV (nA/V)" , fontsize =  'small',c = 'tab:green')
    axs1.tick_params(axis = 'y', colors = 'tab:green')
    plt.tight_layout()

    axs1.annotate('Line Profiles',
                  xy = (0.5, 1+0.2), 
                  xycoords = 'axes fraction',
    horizontalalignment ='center', verticalalignment='top',
    fontsize = 'medium')
    #trim_axs(axs, 2)
    fig2.tight_layout()
    isns.reset_defaults()

    plt.show()
    
    #fig.savefig('fig1.pdf')
    # save first, then show. 
    #plt.show()
    # after plt.show(), figure is reset, 

    return l_pf_start,l_pf_end,fig1,fig2
######################


def line_profile2_xr(xrdata, l_pf_start, l_pf_end, ch_N = [0,2], profile_width = 3):
    """
    
    # Without  GUI input

    Parameters
    ----------
    
    xrdata : xarray  DataSet TYPE
        DESCRIPTION. : input data 
    l_pf_start: np.array([cx,ry])
            astype(int), shape = (2,)
            # cx, ry  idx
    l_pf_end:  np.array([cx,ry])
            astype(int), shape = (2,)
            # cx, ry  idx
    ch_N : Integer list TYPE, optional
            use the channel numbers for line profile
            in the Data Array in the Xarray DataSet
            DESCRIPTION. The default is 0,2.
            the z_fwd_df & LIX_fwd_df
    profile_width : Integer TYPE, optional
            The default is 3.
            line width (perpendicular to the line ) for profile value 
            use the mean value of perpendicular 3 points 

    Returns
    -------

    fig1   : fig object  TYPE
        DESCRIPTION. : figure (line on figures)
    fig2    : fig object TYPE
        DESCRIPTION. : line profile figure (line profiles ) 

    """
    # %matplotlib qt5
    [size_x, size_y] = xrdata.image_size
    
    ch_names_for_lp = []
    for ch_i, ch_name in enumerate(xrdata):
        if ch_i in ch_N : 
            ch_names_for_lp.append(ch_name)
            # check channel names first        
        else :
            pass
    #####################################################################
    l_pf_start  = l_pf_start
    l_pf_end  = l_pf_end
    # given data point
    # %matplotlib inline
    # come back to the inline backend. 
    
    l_pf_length = sp.spatial.distance.euclidean(l_pf_start,l_pf_end)
    l_pf_start_end_pt = [*zip(l_pf_start,l_pf_end)]
    scan_aspect_ratio = (size_x/len(xrdata.X))/(size_y/len(xrdata.Y))
    
    l_pf1 = skimage.measure.profile_line(xrdata[ch_names_for_lp[0]].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
    
    l_pf2 = skimage.measure.profile_line(xrdata[ch_names_for_lp[1]].values, 
                        l_pf_start, 
                        l_pf_end,
                        linewidth = profile_width,
                        reduce_func = np.mean,
                        mode='reflect' )
 
    fig1,axs = plt.subplots (nrows=1,ncols=2, figsize = (6,3))
    isns.imshow(xrdata[ch_names_for_lp[0]].values, 
                robust=True,
                origin ="lower", 
                cmap = 'bone',
                ax = axs[0])
    axs[0].arrow(l_pf_start[0],
             l_pf_start[1],
             l_pf_end[0]-l_pf_start[0],
             l_pf_end[1]-l_pf_start[1],
             width = 3,
             color = 'tab:red')
    axs[0].annotate(ch_names_for_lp[0],
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    if scan_aspect_ratio != 1: 
        fig1.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) 
    # image =  0 :( ==> xrdata[ch_names_for_lp[0]]  )
    isns.imshow(xrdata[ch_names_for_lp[1]].values,
                robust=True,
                origin ="lower",
                cmap = 'bwr',
                ax = axs[1])
    
    axs[1].arrow(l_pf_start[0],
         l_pf_start[1],
         l_pf_end[0]-l_pf_start[0],
         l_pf_end[1]-l_pf_start[1],
         width = 3,
         color = 'tab:green')
    axs[1].annotate(ch_names_for_lp[1],
        xy=(0, 1+0.1), xycoords='axes fraction',
        horizontalalignment='left', verticalalignment='top',
        fontsize='medium')
    if scan_aspect_ratio != 1: 
        fig1.colorbar(isns_channels.get_children()[-2],  
                    fraction = 0.045, 
                    ax = axs[ch_i]) 
    fig1.suptitle(xrdata.title,
                  fontsize = 'small',
                  position=(0.5, 1.0+0.001) )
    fig1.tight_layout()
    isns.reset_defaults()
    plt.show()    
    
    
    
    fig2,axs = plt.subplots (figsize = (4,2))
    axs.plot(np.linspace(0,
                            l_pf_length, 
                            len(l_pf1)
                           )*xrdata.attrs['X_spacing']*1E9,
                l_pf1*1E9,
            c = 'tab:red')
    axs.grid(True)
    plt.xticks(fontsize =  'small')
    plt.xticks(fontsize =  'small')
    axs.set_xlabel( "Length (nm)" , fontsize = 'small')
    axs.tick_params(axis = 'y', colors = 'red')
    axs.set_ylabel( "Height (nm)" , fontsize =  'small',c = 'tab:red')
    
    axs1 = axs.twinx()
    axs1.plot(np.linspace(0,
                        l_pf_length, 
                        len(l_pf1)
                       )*xrdata.attrs['X_spacing']*1E9,
              l_pf2*1E9,
             c = 'tab:green')    
    axs1.set_ylabel( "dI/dV (nA/V)" , fontsize =  'small',c = 'tab:green')
    axs1.tick_params(axis = 'y', colors = 'tab:green')
    plt.tight_layout()

    axs1.annotate('Line Profiles',
                  xy = (0.5, 1+0.2), 
                  xycoords = 'axes fraction',
    horizontalalignment ='center', verticalalignment='top',
    fontsize = 'medium')
    #trim_axs(axs, 2)
    fig2.tight_layout()
    isns.reset_defaults()

    plt.show()
    
    #fig.savefig('fig1.pdf')
    # save first, then show. 
    #plt.show()
    # after plt.show(), figure is reset, 

    return fig1, fig2

##########################

# -




# ### threshold map + labeling area 
#

# +
# select one map & apply thresholds
# choose reference map using bias_mV,
# otsu threholde. 


def th_otsu_roi_label_2D_xr(xr_data, bias_mV_th = 0, threshold_flip = True):
    xr_data_prcssd = xr_data.copy()
    xr_data_prcssd['LDOS_fb_th'] = threshold_otsu_xr (xr_data.sel(bias_mV=bias_mV_th, method="nearest"), threshold_flip= threshold_flip).LDOS_fb
    xr_data_prcssd['LDOS_fb_th_label'] = xr_data_prcssd['LDOS_fb_th'].copy()
    xr_data_prcssd['LDOS_fb_th_label'].values = skimage.measure.label(xr_data_prcssd.LDOS_fb_th.values)
    return xr_data_prcssd
    
    


# +
# select one map & apply thresholds
# choose reference map using bias_mV,
# otsu threholde. 


def th_multiotsu_roi_label_2D_xr(xr_data, bias_mV_th = 200, multiclasses = 3):
    xr_data_prcssd = xr_data.copy()
    xr_data_prcssd['LDOS_fb_th'] = threshold_multiotsu_xr(xr_data.sel(bias_mV=bias_mV_th, method="nearest"), multiclasses = multiclasses).LDOS_fb
    xr_data_prcssd['LDOS_fb_th_label'] = xr_data_prcssd['LDOS_fb_th'].copy()
    xr_data_prcssd['LDOS_fb_th_label'].values = skimage.measure.label(xr_data_prcssd.LDOS_fb_th.values)
    return xr_data_prcssd
    
    


# +
# select one map & apply thresholds
# choose reference map using bias_mV,
# otsu threholde. 


def th_mean_roi_label_2D_xr(xr_data, bias_mV_th = 200, threshold_flip = False):
    '''
    bias_mV_th = threshold_mean_xr at the given bias
    skimage measure  label functino applied to given bias rolling image 
    
    '''
    
    xr_data_prcssd = xr_data.copy()
    xr_data_prcssd['LDOS_fb_th'] = threshold_mean_xr(xr_data.sel(bias_mV=bias_mV_th, method="nearest"), threshold_flip= threshold_flip).LDOS_fb
    xr_data_prcssd['LDOS_fb_th_label'] = xr_data_prcssd['LDOS_fb_th'].copy()
    xr_data_prcssd['LDOS_fb_th_label'].values = skimage.measure.label(xr_data_prcssd.LDOS_fb_th.values)
    return xr_data_prcssd
    
    
    
# -

# #  10. pptx page adding function 
# * **AddAnalysisSlide** function 

# +

def AddAnalysisSlidePage (pptx_file_name,fig_title,fig_file): 
    """
    

    Parameters
    ----------
    pptx_file_name : str TYPE
        DESCRIPTION. pptx file name to add page
    fig_title : str TYPE
        DESCRIPTION. analysis page title 
        use the step-by-step page control for detailed conditions
    fig_file : str TYPE
        DESCRIPTION. figure to add in the page

    Returns
    -------
    None.

    """
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.util import Cm, Pt
    # open previously prepared pptx 
    prs_sxm = Presentation(pptx_file_name)

    sld_slide_layout = prs_sxm.slide_layouts[6] # empty layout 
    
    slide = prs_sxm.slides.add_slide(sld_slide_layout)
    # text box positioning
    left = Inches(0.5)
    top = Inches(0)
    width = Inches(8)
    height = Inches(0.5)

    slide_textbox = slide.shapes.add_textbox(left, 
                                               top, 
                                               width, 
                                               height) # add text box 


    paragraph0 = slide_textbox.text_frame.add_paragraph()
    # To control the font, use the paragraph
    paragraph0.text = fig_title # (그림제목) 내용설정
    paragraph0.level = 1   
    paragraph0.font.size = Pt(16)
    # adjust the font size
    
    ###########
    # figure selection 
    # check the dpi for the saved file
    
    # check the imaage positions in page 
    left = Inches(0.5)
    top = Inches(1)

    figure_img = slide.shapes.add_picture(fig_file,
                                           left, 
                                           top)
    
    prs_sxm.save(pptx_file_name)  

    return

    

# -

# # Interactive plot
#



#
# * xr_isns_plot
#     * xr_data channels isns plot. 
#
#
#
#
# * xr_isns_plot_r_space
#     *   real space (r) XR data set plot (xr_data) 
#         * input xrdata, 
#         * use seaborn image ( origin = "lower") 
#         * **xr_isns_plot_r_space** (xr_data, ncols = 2, figsize = (5,5)) $\rightarrow$  return fig
#   
# >* (e.g.) xr_data_set_isns_plot_r_space(z_LIX_fNb_xr)
#
# * xr_isns_plot_k_space
#     * momentum(k) space XR data set plot (xr_data_fft)
#         * **xr_isns_plot_k_spac** (xr_data_fft, ncols = 4, figsize =  (8,4))$\rightarrow$    return fig
#
# >* (e.g.) xr_data_set_isns_plot_k_space(z_LIX_fNb_xr_fft)

# # skimage exposure
#

# + id="LtsdRqHwVAu2" jp-MarkdownHeadingCollapsed=true
#######################################
# plot Xarray (xr) data set with isns
###############\###############\
# len(z_LIX_fNb_xr_intrp.data_vars) # check xr DataArray number
 #######################
def xr_isns_plot_r_space(xrdata,
                         ncols = 2,
                         figsize = (6,6)): 
    """
    display xr DATA Array channels 
    use isns. 'nm' unit, 
    Parameters
    ----------
    xrdata : Xarray DataSet TYPE
        DESCRIPTION.
    ncols : Integer TYPE, optional
        DESCRIPTION. The default is 2.
        number of columns to grid image display
    figsize : tuple TYPE, optional
        DESCRIPTION. The default is (6,6).
        out put figure size 

    Returns
    -------
    fig :image show TYPE
        DESCRIPTION.
        need to chekc the figure saving 
        
    """
    # aspect ratio check 
    [size_x, size_y] = xrdata.image_size
    step_dx = size_x/len(xrdata.X)
    step_dy = size_y/len(xrdata.Y)
    scan_aspect_ratio = (size_x/len(xrdata.X))/(size_y/len(xrdata.Y))
    if scan_aspect_ratio == 1: 
        cbar_show = True
    else:
        cbar_show = False
    ####################################
    isns.set_image(cmap='viridis', origin ="lower")
    fig,axes = plt.subplots(ncols = ncols,
                            nrows = len(xrdata.data_vars)//ncols+1,
                            figsize = figsize)
    axs = axes.ravel()
    axs = trim_axs(axs, len(xrdata.data_vars))
    # triming first with respect to # of data channel 
    isns.set_image(cmap = 'viridis', origin='lower')   #  set image  direction
    
    for i,i_channels in enumerate(xrdata):
        isns_channels= i_channels+'_isns'
        print(isns_channels)
        if 'z_' in  i_channels:
            cmap =  'copper'
        elif 'Z_' in  i_channels:
            cmap =  'copper'
        elif 'LIX_' in  i_channels:
            cmap = 'bwr'
        isns_channels = isns.imshow(xrdata[i_channels],
                                    ax =  axs[i],
                                    cbar = cbar_show,
                                    dx = step_dx*1E9,
                                    units = "nm",
                                    cmap = cmap,
                                    robust = True)
        axs[i].set_title(i_channels, loc='left', fontsize = 'small')
        #################################
        if scan_aspect_ratio != 1: 
                fig.colorbar(isns_channels.get_children()[-2],  
                            fraction = 0.045, 
                            ax = axs[i])  
        ################################
        # colorbar setting separately 
        ############################## 
        if 'ref_lattice_a0' in xrdata.attrs:
            ref_lattice_r0 = xrdata.attrs['ref_lattice_a0'] # kspace atomic lattice_k0
            ref_6p_radious = ref_lattice_r0/xrdata.attrs['X_spacing']
        else: pass

    # 4channel 2x2 
    fig.suptitle(xrdata.title, fontsize = 'small', position=(0.5, 1.0+0.05) )
    fig.tight_layout()
    isns.reset_defaults()
    plt.show()
    return fig
# test
# xrdata_set_isns_plot_r_space(z_LIX_fNb_xr_intrp)


#####################################################################
def xr_isns_plot_k_space(xrdata_fft,
                         nrows = 2, 
                         ncols = 2,
                         figsize =  (8,4), 
                         zoom_in_fft = True,
                         zoom_in_expand = 3): 
    """
    display xr DATA Array channels after FFT 
    (k-space)
    default cmap = 'viridis'
    # use isns. 'nm' unit, 
    Parameters
    ----------
    xrdata_fft : Xarray DataSet TYPE
        DESCRIPTION.
        xrdata_fft  for k space plot 
    ncols : Integer TYPE, optional
        DESCRIPTION. The default is 4.
        number of columns to grid image display
    figsize : tuple TYPE, optional
        DESCRIPTION. The default is (8,4).
        out put figure size
    zoom_in_fft : Boolean TYPE, optional
        DESCRIPTION. The default is True.
        check the zoomin or not 
    zoom_in_expand :  Integer TYPE, optional
        DESCRIPTION. The default is 3.
        Zoom in area : ref_lattice_k0 * zoom_in_expand

    Returns
    -------
    fig : out figure TYPE
        DESCRIPTION.
        check the saving figure results. 
        

    """
    # aspect ratio check 
    [size_x, size_y] = xrdata_fft.image_size
    step_dx = size_x/len(xrdata_fft.freq_X)
    step_dy = size_y/len(xrdata_fft.freq_Y)
    scan_aspect_ratio = (size_x/len(xrdata_fft.freq_X))/(size_y/len(xrdata_fft.freq_Y))
    if scan_aspect_ratio == 1: 
        cbar_show = True
    else:
        cbar_show = False
    #######################################
    isns.set_image(cmap='cividis', origin ="lower")
    
    
    #  set image  direction
    fig,axes = plt.subplots(ncols = ncols,
                            nrows = len(xrdata_fft.data_vars)//ncols+1,
                            figsize = figsize)
    axs = axes.ravel()
    axs = trim_axs(axs, len(xrdata_fft.data_vars))
    # triming first with respect to # of data channel 
    for i,i_channels in enumerate(xrdata_fft):
        if 'z_' in i_channels:
            cmap = 'Greys'
        elif 'Z_' in i_channels:
            cmap = 'Greys'
        elif 'LIX_' in i_channels:
            cmap = 'Blues'
        
        isns_channels= i_channels+'_isns'
        if i_channels.endswith('fft'): # fft channels 
            dx=(1/size_x)*1E-9
            units="1/nm"
            dimension  = "si-reciprocal"
        else : 
            dx=step_dx*1E9
            units="nm"
            dimension  = "si"
        print(isns_channels)

        isns_channels = isns.imshow(xrdata_fft[i_channels],
                                    ax =  axs[i],
                                    cmap= cmap,
                                    cbar = cbar_show,
                                    dx = dx,
                                    units= units,
                                    dimension = dimension,
                                    robust = True)
        axs[i].set_title(i_channels, loc='left', fontsize = 'small')
        #################################
        if scan_aspect_ratio != 1: 
            fig.colorbar(isns_channels.get_children()[-2],  
                        fraction = 0.045, 
                        ax=axs[i])  
        ################################
        # colorbar setting separately
        ################################
        if 'ref_lattice_k0' in xrdata_fft.attrs: 
            ref_lattice_k0 = xrdata_fft.attrs['ref_lattice_k0'] # kspace atomic lattice_k0
            ref_6p_radious = ref_lattice_k0/xrdata_fft.attrs['freq_X_spacing']
        else : pass
    
        ## for zoom in (
        if zoom_in_fft == True:
            # center position (0,0) -> (0,0) in (freq_X,freq_Y )-> find index
            # np.abs(z_LIX_fNb_xr_fft.freq_X).argmin()
            x_lower_limit = \
             np.abs(xrdata_fft.freq_X).argmin()\
             -int(ref_6p_radious)*zoom_in_expand
            x_upper_limit = \
            np.abs(xrdata_fft.freq_X).argmin()\
            +int(ref_6p_radious)*zoom_in_expand 
            axs[i].set_xlim([x_lower_limit,x_upper_limit])
            axs[i].set_ylim([x_lower_limit,x_upper_limit])    
        ## for zoom in )
    fig.suptitle(xrdata_fft.title +'(2D FFT)',
                 fontsize = 'small',
                 position=(0.5, 1.0+0.05) )
    fig.tight_layout()
    isns.reset_defaults()
    plt.show()
    
    return fig


# # + [markdown] tags=[] jp-MarkdownHeadingCollapsed=true
# ### 9.  xrdata_fft analysis 
# * (p6 & r6,rot_angle, pad & rot)
# * for xrdata_fft find 6 points 
# * 1-2-7 : xr_data analysis ( p6r6, rot angle, padding, rotate, ) 
#
# ### find 6 peaks and 6 reference points 
# * xrdata_fft data  should have <U> reference$a_{0}$ </U>
# * **xrdata_fft_plot_p6r6** (xrdata_fft,zoom_in_fft = True, zoom_in_expand = 3):
# > * input : xrdata_fft 
# > * output: p6,r6,fig
# * p6,r6 point lists are saved as attributes to the xr_data_ffft data 

# # + tags=[]
##################################################################
# input data :  xrdata_fft 
#~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
# return #  peak6points_kspace_idxs, ref_6points_kspace_idxs
#~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
#################################

def xrdata_fft_plot_p6r6(xrdata_fft,
                           zoom_in_fft = True, 
                           zoom_in_expand = 3):
    '''
    # xrdata_fft should have attributes 'ref_lattice_a0' & 'ref_lattice_k0'
    # xrdata_fft  # with ref_a0  : input data     
    # zoom_in_fft : zoomin for the center region 
    # zoom_in_expand : control zoom in area (3times of atomic lattice)
        out put  p6 & r6 
            p6 : peak_6points_kspace_idxs,
            r6: ref_6points_kspace_idxs
    
    #~~~~~~#
    '''
    ref_lattice_a0 = xrdata_fft.attrs['ref_lattice_a0'] # ref_lattice_a0
    ref_lattice_k0 = xrdata_fft.attrs['ref_lattice_k0'] # kspace atomic lattice_k0
    ref_6pts  = ref_lattice_k0 * np.array([[math.cos(pt_i* math.pi/3), 
                                   math.sin(pt_i* math.pi/3)]
                       for pt_i in range(6)])
    # ref_6pts: (x,y) form
    # from a0 calibrate  'k0' & point list
    ref_6points_kspace_x =  (ref_6pts[:,0] 
                            - xrdata_fft.freq_X.values[0]
                            )/xrdata_fft.freq_X.spacing
    # 
    ref_6points_kspace_y =  (ref_6pts[:,1] 
                            - xrdata_fft.freq_Y.values[0]
                            )/xrdata_fft.freq_Y.spacing
    ref_6p_radious = ref_lattice_k0/xrdata_fft.freq_X.spacing
    # fro "skimage.draw.disk" center_indx -> dtype(int), save as np 
    ref_6points_kspace_idxs = np.array([[round(x,2),round(y,2)]
            for x, y  in zip(ref_6points_kspace_x, ref_6points_kspace_y) ])
    # use list comprehention x,y 
    # round ( , 2)
    #print(ref_6points_kspace_idxs)
    # ref_6points_kspace_idxs : (x,y) form
    #~~~~~~~~~~~~~~~~~~~~~~~~~#

    # find  6 point  in the measured fft 
    peak_6points_kspace_idxs = np.array(
    [np.zeros_like(ref_6points_kspace_idxs) 
    for i in range (len(xrdata_fft))])
    p6s_k_pts = peak_6points_kspace_idxs.copy()
    # assigne array for the measured peak lists (for each channels)
   
    # ref6points_kspace_idxs : 
    # ideal 6 points of atomic lattice
    # peak6points_kspace_idxs :
    # measured 6 points from atomic lattice
    # ref point. shape = (6,2) 
    # local peaks. shape  = (N, 6, 2) dimension 
    # 
    ################### plot #################
    isns.set_image( origin = "lower", cmap = 'cividis')
    
    fig,axes = plt.subplots(2,2, figsize = (6,6))
    axs = axes.ravel()
    for i,ch_map in enumerate (xrdata_fft):
        isns_channel = ch_map+'isns'
        print(ch_map)    
        # use the isns instead of xr.plot 
        # plot type is different, isns type is more controllable
        # for example, we can adjust the set_aspect ration
        isns_channel = isns.imshow(xrdata_fft[ch_map],
                                   ax = axs[i],
                                   dx = xrdata_fft.freq_X.spacing*1E-9,
                                   units= "1/nm",
                                   dimension = "si-reciprocal",
                                   robust = True)
        axs[i].set_aspect(1)
        
        #################################
        # Red circle for peak positions  plot
        #################################
        for r_i, ref_6points_kspace_idx in  enumerate(ref_6points_kspace_idxs):
            ref_circles = Circle(tuple(ref_6points_kspace_idx), # center x,yidxs
                                    ref_6p_radious/4, # radius
                                    facecolor = 'none', 
                                    edgecolor = 'red',
                                    linewidth = 1,
                                    alpha = 0.5)
            axs[i].annotate(str(r_i),
                            tuple(ref_6points_kspace_idx),
                            color='r',
                            size =8) # ref peak numbers
            axs[i].add_patch(ref_circles)
            # circle xy = (x,y) 
            #  ref_6points_kspace_idx is also  (x, y) 
        #################################
        # white circles for local peaks
        # find local peaks (it could be any number)
        #################################
        local_peaks_in_ch_rycx_test = skimage.feature.peak_local_max(
            gaussian_filter(
                xrdata_fft[ch_map].values, sigma=1
                ),
                min_distance=int(ref_6p_radious/2),
                threshold_rel = 0.2 )
        # check the number of local peaks 
        # if there is no local peaks, adjust the peak_local_max prameters 
        print ('local leaks Number : ', len(local_peaks_in_ch_rycx_test))
        # if "local leaks Numbe" is zero ==> error!
        if len(local_peaks_in_ch_rycx_test) == 0: 
            local_peaks_in_ch_rycx = skimage.feature.peak_local_max(
                gaussian_filter(xrdata_fft[ch_map].fillna(0).values,
                                sigma=1),
                min_distance=int(ref_6p_radious/2) 
            )
                
            # Find all local peaks without threshold condition
            #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
            # in case of filtered FFT --> fill the nan to 0 
            #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
            # local_peaks_in_ch result ==> (ry,cx) order due to 'peak_local_max'
            local_peaks_in_ch = local_peaks_in_ch_rycx[:,::-1]
        # change the array (ry, cx,) -> (cx,ry) 
            print ('local leaks Number : ', len(local_peaks_in_ch))
        else :
            local_peaks_in_ch_rycx = skimage.feature.peak_local_max(
            gaussian_filter(
                xrdata_fft[ch_map].values, sigma = 1
                ),
                min_distance=int(ref_6p_radious/2),
                threshold_rel = 0.2 )
            # Find local loca peaks with top 20% in the (channel) image 
            # local_peaks_in_ch result ==> (ry,cx) order due to 'peak_local_max'
        local_peaks_in_ch = local_peaks_in_ch_rycx[:,::-1]
            

        #################
        for local_peak in  local_peaks_in_ch:
            #print(local_peak)
            peak_bumps = Circle(tuple(local_peak),
                                ref_6p_radious/3, # radius
                                facecolor = 'none', 
                                edgecolor = 'white',
                                linewidth = 1,
                                alpha = 0.5) 
            axs[i].add_patch(peak_bumps)
            # circle xy = (x,y) 
            #  local_peak  is also  (x, y) 
            # white circles. 
        ''' # if you want use scatter
        axs[i].scatter(local_peaks_in_ch[:,0], # cx
                       local_peaks_in_ch[:,1], # ry
                       c='w', 
                       s = ref_6p_radious/2,
                       alpha =0.5
                       )'''
        #######################################
        # select (slct) 6 local peaks 
        # which is close to the ref 6 pts 
        ###########################################
        find_6peaks_idx_in_peak_list = [
                                        np.array([distance.euclidean(
                                            ref_6points_kspace_idx,
                                            local_peak
                                            ) 
                                        for local_peak in local_peaks_in_ch]).argmin()
                                        for ref_6points_kspace_idx 
                                        in  ref_6points_kspace_idxs
                                        ]
        print (find_6peaks_idx_in_peak_list)
        # local peaks in  each channels 
        # make list for euclidean distance with one of ref pt
        # check (x,y)  vs (x,y) !!!
        # find 'arg min '
        # 
        peak_6points_kspace_idx = local_peaks_in_ch[find_6peaks_idx_in_peak_list]
        
        p6_k_pts_cx = xrdata_fft.freq_X[peak_6points_kspace_idx[:,0].astype(int)]
        
        p6_k_pts_ry = xrdata_fft.freq_Y[peak_6points_kspace_idx[:,1].astype(int)]
        p6_k_pts_i = np.stack((p6_k_pts_cx,p6_k_pts_ry),axis =1)
        # index --> real k-point values w.r.t coordinates
        # for [i]th channel 
        
        p6s_k_pts[i] = p6_k_pts_i
        peak_6points_kspace_idxs[i] = peak_6points_kspace_idx

        ############################################
        # add selected local  6 peaks 
        ############################################
        axs[i].scatter(peak_6points_kspace_idx[:,0],
                       peak_6points_kspace_idx[:,1],
                       c='b',
                       s = ref_6p_radious/1,
                       alpha =0.5)
        for pk_j, peak6points in  enumerate(peak_6points_kspace_idx):
            axs[i].annotate(str(pk_j),
                            xy = tuple(peak6points),
                            color='b',
                            size =8) #anntates
        # peak_6points_kspace_idxs .shape = [i]channel * (6,2)
        # check channel [i]
        #################
        '''
        for pk_i, local_peak_in6 in  enumerate(peak_6points_kspace_idxs[i]):
            local_peak_6list = Circle(tuple(local_peak_in6),
                                ref_6p_radious/4, # radius
                                facecolor = 'none', 
                                edgecolor = 'blue',
                                linewidth = 1,
                                alpha = 0.8) 
            axs[i].add_patch(local_peak_6list)
            axs[i].annotate(str(pk_i),
                            tuple(local_peak_in6)
                            , color='blue',
                            size =8) #번호넣기
                            
                            '''
        ######################################
        #Zoom in with  xy limit
        if zoom_in_fft == True:
            # center position (0,0) -> (0,0) in (freq_X,freq_Y )-> find index
            # np.abs(z_LIX_fNb_xr_fft.freq_X).argmin()
            x_lower_limit = \
             np.abs(xrdata_fft.freq_X).argmin()\
             -int(ref_6p_radious)*zoom_in_expand
            x_upper_limit = \
            np.abs(xrdata_fft.freq_Y).argmin()\
            +int(ref_6p_radious)*zoom_in_expand 
            axs[i].set_xlim([x_lower_limit,x_upper_limit])
            axs[i].set_ylim([x_lower_limit,x_upper_limit])
    isns.reset_defaults()
    plt.tight_layout()
    
    plt.show()
    
    
    # assign r6 & p6s #xrdata_fft.attributes  
    xrdata_fft.attrs['r6_k_pts'] = ref_6pts
    xrdata_fft.attrs['p6s_k_pts'] = p6s_k_pts
    
    #print ('6 reference spots(dst)')
    #print (ref_6points_kspace_idxs)
    #print ('6 measured peaks(src)')
    #print (peak_6points_kspace_idxs)
    return  peak_6points_kspace_idxs, ref_6points_kspace_idxs, fig
# return p6 & r6 : check the order 



# # + colab={"base_uri": "https://localhost:8080/"} id="mGfdySXL55C8" outputId="d2e5f006-392c-4ffa-a70f-f08ab139b707"
###################################
# src & dst points 
#dst_r6 = z_LIX_fNb_xr_fft_r6
#src_p6 = z_LIX_fNb_xr_fft_p6

def xrdata_fft_rot_angle(xrdata_fft): 
    '''
    only after : xrdata_fft_plot_p6r6
    input: xrdata_fft
    output : rotation_angle_deg (deg)
    after get rotation_angle_deg, 
    send the xrdata_fft['rot_angle']  to the (real space) xrdata
    eg) #z_LIX_fNb_xr_lattice.attrs = z_LIX_fNb_xr_lattice_fft.attrs

    '''
    # p6 & r6 from attributes
    src_p6 = xrdata_fft.attrs['p6s_k_pts']
    dst_r6 = xrdata_fft.attrs['r6_k_pts']
    
    ###################################
    # check the slope  (dy/dx)
    #####################################
    # use the horizontal axis  spots (0, 3) 
    dst_slope = np.degrees(
            2*math.pi+math.atan2(
                (dst_r6[0]-dst_r6[3])[1],
                (dst_r6[0]-dst_r6[3])[0]
                )
            )
    #####################################
    # measured p6points 
    # check the different dimension 
    # use the first channel (z_fwd)
    #####################################
    src_slope = np.degrees(
            2*math.pi+math.atan2(
                (src_p6[0][0]-src_p6[0][3])[1],
                (src_p6[0][0]-src_p6[0][3])[0]
                )
            )
    rotation_angle_deg =  src_slope - dst_slope 
    xrdata_fft.attrs['rot_angle'] = rotation_angle_deg
    return rotation_angle_deg


################################
#test 
#rotation_angle_deg =   xrdata_fft_rot_angle(z_LIX_fNb_xr_lattice_fft)
#
# after find rotational angle, copy the attributes to the real space data 
#z_LIX_fNb_xr_lattice.attrs = z_LIX_fNb_xr_lattice_fft.attrs

#print(rotation_angle_deg)
#print (z_LIX_fNb_xr_lattice_fft)

# # + [markdown] tags=[]
# * fig can be saved as image (png) 
#
# ### find_rot_angle from (p6, r6)
#
# * **xrdata_fft_rot_angle**
# > * input : xrdata_fft
# > * output : rot_angle_deg
#
#     * based on previous 6 spot test 
#         * find slope between reference 6points & measured data
#     * xr padding region checking (pad_x, pad_y)
#     * Rotation angle is saved as attributes 
#         * save it to both xrdata_fft & xr_data
#         
#
# ### Padding according to (rot_angle) 
# * **xrdata_padding**  after FFT 
#     * Padding original data to prevent image edge cutting off during affine transform 
#     * rotating the original data set 
#         * rotation test (resize (True) during rotation) 
#         * Confirm the required area for **xarray padding**
#     * adjust **x,y coords**  according to the padding size
#         *re-calibrate: spacing + offset
#         
# ### Rotating image after (padding) 
# * **xrdata_rotate** after padding
#     * rotate the original xrdata after padding
#     * increased size after rotation <-- padding size
# -

# ### 7.  FFT masking 
# * fft_masking_center 
# * data_fft_cmplx_iff

# +
def fft_masking_center (xrdata_fft_cmplx,
                        ch_N_to_show = 0,
                        pt_input = 3,
                        zoom_in_fft = True,
                        zoom_in_expand = 2):
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
    # input data (xrdata_fft) should be complex 
    # use the 'complex_output = True' option for  twoD_fft_xr
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~#
    ch_name_list = [ch_name for ch_name in xrdata_fft_cmplx ]
    print('channel to show : ' , ch_name_list[ch_N_to_show])
    #list for ch_name
    fft_ch_to_show = xrdata_fft_cmplx[ch_name_list[ch_N_to_show]].values
    #print (fft_ch_to_show)
    
    if fft_ch_to_show.dtype == 'complex128':
        # masking centeral area with average distance 
        # change the matplotlib backend : inline --> qt5  
        # Open a new window 
        # %matplotlib qt5
        fig,axs = plt.subplots (nrows=1, ncols=1, figsize = (5,5))
        isns.imshow(np.abs(fft_ch_to_show),
                    ax = axs,
                    robust=True,
                    origin ="lower")
        gui_pts = fig.ginput(pt_input)
        print(gui_pts)
        gui_pts_idx = np.array([*gui_pts]).astype(int)# dont forget '*'
        print(gui_pts_idx) # points index from GUI
        plt.show()
    
        # %matplotlib inline
        # comeback to inline backend plot.

        xrdata_fft_cmplx_masked = xrdata_fft_cmplx.copy()
        #prepare the mask 
        fft_fltrng_mask = np.zeros (fft_ch_to_show.shape)
        center_rycx = np.array([*fft_fltrng_mask.shape])//2
        fltrng_mask_radious = np.array(
            [ distance.euclidean(center_rycx , gui_pt) 
             for gui_pt in gui_pts_idx]).mean() # averaged distances (center -GUI points)
        print('distance_from_center : ', fltrng_mask_radious)

        fft_fltrng_mask_ry,fft_fltrng_mask_cx = skimage.draw.disk(
            center = (center_rycx[0],center_rycx[1]),
            radius = fltrng_mask_radious)
        fft_fltrng_mask[fft_fltrng_mask_ry,
                        fft_fltrng_mask_cx ] = 1
        
        # prepare copied channel   
        xrdata_fft_cmplx_masked['fft_fltrng_mask'] = xrdata_fft_cmplx_masked[ch_name_list[ch_N_to_show]]
        # replace the masked 
        xrdata_fft_cmplx_masked['fft_fltrng_mask'].values = fft_fltrng_mask
        # we can remove the mask channel later 
        for ch_name in xrdata_fft_cmplx_masked: 
            xrdata_fft_cmplx_masked[ch_name] = xrdata_fft_cmplx_masked[ch_name].where(
            xrdata_fft_cmplx_masked['fft_fltrng_mask']==1)   
        xrdata_fft_cmplx_masked = xrdata_fft_cmplx_masked.drop(["fft_fltrng_mask"])
        
        ############## 
        # plot masked     

    else: print ('Error: it is not complex128')
    return xrdata_fft_cmplx_masked


#test
# xrdata_fft_cmplx_masked = fft_masking_center(xrdata_fft_cmplx)
#

#complex_input & ifft & save absolute value 

def xrdata_fft_cmplx_iff(xrdata_fft_cmplx): 
    ch_name_list = [ch_name for ch_name in xrdata_fft_cmplx ]
    #xrdata_fft_cmplx[ch_name_list[0]]
    if xrdata_fft_cmplx[ch_name_list[0]].values.dtype != 'complex128':
        print ('Error: it is not complex128')
    else :
        xrdata_fft_cmplx_ifft = xrdata_fft_cmplx.copy()
        for ch_name in xrdata_fft_cmplx: 
            print(ch_name[:-4])
            xrdata_fft_cmplx_ifft.rename({ch_name: ch_name[:-4]}) # change channel names
            xrdata_fft_cmplx_ifft[ch_name[:-4]] = xrft.xrft.ifft(xrdata_fft_cmplx[ch_name].fillna(0))
            # ifft & fill nan --> 0 for ifft 
            xrdata_fft_cmplx_ifft[ch_name[:-4]].values = np.abs( xrdata_fft_cmplx_ifft[ch_name[:-4]])
            # make an absolute after ifft 
            xrdata_fft_cmplx_ifft.drop([ch_name])
            #xr_data_after_fft= xr_data_after_fft.drop_dims(["X", "Y"])
        xrdata_fft_cmplx_ifft = xrdata_fft_cmplx_ifft.drop_dims(["freq_X","freq_Y"])
    return xrdata_fft_cmplx_ifft
                                  
##  test 
## xrdata_fft_cmplx_iff_temp =  xrdata_fft_cmplx_iff(xrdata_fft_cmplx_masked)




# successfully masked 
# -


# ### Xr rotation function 
#

# ### 2D rotation 

### Xr rotation function 
# rotate the XY plan in xr data 
def rotate_2D_xr (xrdata, rotation_angle): 
    # padding first 
    for ch_i,ch_name in enumerate (xrdata):
            if ch_i == 0:  # use only the first channel to calculate a padding size 
                padding_shape = skimage.transform.rotate(xrdata[ch_name].values.astype('float64'),
                                                         rotation_angle,
                                                         resize = True).shape          


                padding_xy = (np.array( padding_shape)-np.array(xrdata[ch_name].shape) +1)/2
                padding_xy = padding_xy.astype(int)
    xrdata_pad = xrdata.pad(X=(padding_xy[0],padding_xy[0]), 
                            Y =(padding_xy[1],padding_xy[1]),
                            mode='constant',
                            cval = xrdata.min())
    if np.array(xrdata_pad[ch_name]).shape != padding_shape:
        # in case of xrdata_pad shape is +1 larger than real padding_shape
        # index 다루는 법  (X)
        x_spacing = np.diff(xrdata.X).mean()
        y_spacing = np.diff(xrdata.Y).mean()
        xrdata.X[0]
        xrdata.Y[0]

        x_pad_dim = padding_shape[0]#int(padding_xy[0]*2+xrdata.X.shape[0])
        y_pad_dim = padding_shape[1]#int(padding_xy[0]*2+xrdata.Y.shape[0])

        x_pad_arr =  np.linspace(-1*padding_xy[0]*x_spacing, x_spacing*x_pad_dim,x_pad_dim+1)
        y_pad_arr =  np.linspace(-1*padding_xy[1]*y_spacing, y_spacing*y_pad_dim,y_pad_dim+1)

        # 0 에서 전체 크기 만큼 padding 한결과를 array 만들고 offset 은 pad_x 만큼 
        x_pad_arr.shape
        y_pad_arr.shape
        xrdata_pad = xrdata_pad.assign_coords( {"X" :  x_pad_arr}).assign_coords({"Y" :  y_pad_arr})
        xrdata_rot = xrdata_pad.sel(X = xrdata_pad.X[:-1].values, Y = xrdata_pad.Y[:-1].values)
        print ('padding size != rot_size')

    else : # np.array(xrdata_pad[ch_name]).shape == padding_shape 
            # in case of xrdata_pad shape is +1 larger than real padding_shape

        # index 다루는 법  (X)
        x_spacing = np.diff(xrdata.X).mean()
        y_spacing = np.diff(xrdata.Y).mean()
        xrdata.X[0]
        xrdata.Y[0]

        x_pad_dim = padding_shape[0]#int(padding_xy[0]*2+xrdata.X.shape[0])
        y_pad_dim = padding_shape[1]#int(padding_xy[0]*2+xrdata.Y.shape[0])

        x_pad_arr =  np.linspace(-1*padding_xy[0]*x_spacing, x_spacing*x_pad_dim,x_pad_dim)
        y_pad_arr =  np.linspace(-1*padding_xy[1]*y_spacing, y_spacing*y_pad_dim,y_pad_dim)

        # 0 에서 전체 크기 만큼 padding 한결과를 array 만들고 offset 은 pad_x 만큼 
        x_pad_arr.shape
        y_pad_arr.shape
        xrdata_pad = xrdata_pad.assign_coords( {"X" :  x_pad_arr}).assign_coords({"Y" :  y_pad_arr})
        xrdata_rot = xrdata_pad.copy()
        print ('padding size == rot_size')

    for ch in xrdata:
            xrdata_rot[ch].values = skimage.transform.rotate(xrdata[ch].values.astype('float64'),
                                                             rotation_angle,
                                                             resize = True,
                                                             cval = xrdata[ch].values.astype('float64').min())
        # after rotation, replace the padding to rotated image 

    # rotate with ( resize= True ) 
    # after rotation, replace the padding to rotated image 
    return xrdata_rot


# # # More?
#
# --
# ---

# *  drift compensation function. .
# *  conceptually right. 
#     * working performance is not good yet..)
# * not perfectly working. 
# * need to improve ... (for loop...) ==>def another function  to use it in the list comprehesion  use? 
# * apply the same mechnism to the 2D channels 
# * make a 3D function 
#
# * 

def drift_compensation_y_topo_crrltn (xr_data_topo, y_sub_n=5, drift_interpl_method='nearest'): 
    y_N = len (xr_data_topo.Y)
    y_sub_n = y_sub_n
    #y_j = 0 
    offset = np.array([0, y_N//2])
    # use for loop 
    print ('only for topo, 2D data, apply to 3D data & other channels later ')
    for y_j  in range (len (xr_data_topo.Y)//y_sub_n - 1) :
        y_N = len (xr_data_topo.Y)
        #print (y_j)

        Y_sub_n0 = y_j*y_sub_n * xr_data_topo.Y_spacing
        Y_sub_n1 = (y_j+1)*y_sub_n * xr_data_topo.Y_spacing
        Y_sub_n2 = (y_j+2)*y_sub_n * xr_data_topo.Y_spacing
        #print (Y_sub_n0, Y_sub_n1, Y_sub_n2)
        # check Y drift comparision area 
        # use y_sub_n = 5 ==> 0-5, 6-10, 10-5, ... 
        line0 = xr_data_topo.where(xr_data_topo.Y >= Y_sub_n0, drop = True).where (xr_data_topo.Y < Y_sub_n1, drop = True ).topography
        line1 = xr_data_topo.where(xr_data_topo.Y >=  Y_sub_n1, drop = True).where (xr_data_topo.Y <  Y_sub_n2, drop = True ).topography
        # select two region for correlation search 
        corrl_line0_line1 = sp.signal.correlate2d(line0.values, line1.values, mode = 'same')#  use mode = same area to use the line0 X& Y value
        # search for the correlation. if correlation is not center --> drift. 
        # but.. there will be an step edge (horizontal), or atomic lattice --> y_sub_n << atomic lattice 
        ind_max = np.array (np.unravel_index(np.argmax(corrl_line0_line1, axis=None), corrl_line0_line1.shape)) # find max point 
        # find argmax index point
        #print (ind_max)
        offset = np.vstack ([offset, ind_max])
    
    offset_0 = offset[: , -1] -  y_N//2
    # check offset from center 
    #offset_accumulation  = [ offset_0[:n+1].sum()  for n in range (len(offset_0)) ]
    
    offset_accumulation  = np.array ( [ offset_0[:n].sum()  
                                       for n in range (len(offset_0)+1) ])*xr_data_topo.Y_spacing 
    # offset is from between two region.. get an accumlated offset. for whole Y axis. 
    offset_accumulation_df =pd.DataFrame (
        np.vstack ([ np.array ([ y_j *y_sub_n *xr_data_topo.Y_spacing  
                                for y_j in range(len (xr_data_topo.Y)//y_sub_n+1) ]), 
                    offset_accumulation]).T, columns  =['Y','offset_X'])
    offset_accumulation_xr  = offset_accumulation_df.set_index('Y').to_xarray()
    offset_accumulation_xr_intrpl = offset_accumulation_xr.offset_X.interp(Y = xr_data_topo.Y.values,  method=drift_interpl_method)
    # accumluted offset--> covert to df , xr, 
    # accumnulated offset to compensate in X 
    # use interpolation along Y --> point offset calc ==> apply to all y points. 

    # for each lines, adjust value after offset compensated  ==> interpolated again. 
    xr_data_topo_offset = xr_data_topo.copy(deep= True)
    # dont forget deep copy... 
    
    for y_j, y  in enumerate (xr_data_topo.Y):
        new_x_i =  xr_data_topo.isel (Y=y_j).X - offset_accumulation_xr_intrpl.isel(Y=y_j)
        # for each y axis. shift X position 
        xr_data_topo_offset_y_j = xr_data_topo_offset.isel (Y=y_j).assign_coords({"X": new_x_i}) 
        # assign_coord as a new calibrated offset-X coords
        xr_data_topo_offset_y_j_intp = xr_data_topo_offset_y_j.interp(X=xr_data_topo.X)
        # using original X points, interpolated offset- topo --> set new topo value to original X position 
        xr_data_topo_offset.topography[dict(Y = y_j)] =  xr_data_topo_offset_y_j_intp.topography
        #grid_topo_offset.isel(Y=y_j).topography.values = grid_topo_offest_y_j_intp.topography
        # use [dict()] for assign values , instead of isel() 
        # isel is not working... follow the instruction manual in web.!
    fig,axs = plt.subplots(ncols = 2, figsize = (6,3))
    xr_data_topo.topography.plot(ax =axs[0])
    xr_data_topo_offset.topography.plot(ax =axs[1])
    plt.show()
    return xr_data_topo_offset



# +
# kde plot gaussian fitting multi peak cases 
# code from bing chat 

# +
#######################################
## Generated code from the bing chat!
########################################

# input_pd : 2 column pandas dataframe 


def multi4gaussian_fit_pd (pd_df, guess = [1E9, -0.3E-10, 1E-10,
         1E9, -0.1E-10, 1E-10, 
         1E9, 0.1E-10, 1E-10,
         1E9, 0.4E-10, 1E-10, 
         0]):
    '''
    # input_pd : 2 column pandas dataframe 
    # guess  :
    # Initial guesses for the parameters to fit: 4 amplitudes, means and standard deviations plus a continuum offset.
    '''
    # gauss



    from scipy.optimize import curve_fit
    import matplotlib.pyplot as plt
    import numpy as np

    def gaussian(x, A, x0, sig):
        return A*np.exp(-(x-x0)**2/(2*sig**2))

    def multi_gaussian(x, *pars):
        offset = pars[-1]
        g1 = gaussian(x, pars[0], pars[1], pars[2])
        g2 = gaussian(x, pars[3], pars[4], pars[5])
        g3 = gaussian(x, pars[6], pars[7], pars[8])
        g4 = gaussian(x, pars[9], pars[10], pars[11])
        return g1 + g2 + g3 + g4 + offset

    vel, flux = pd_df.iloc[:,0],  pd_df.iloc[:,1]

    # Initial guesses for the parameters to fit: 4 amplitudes, means and standard deviations plus a continuum offset.
    guess = [1E9, -0.3E-10, 1E-10,
             1E9, -0.1E-10, 1E-10, 
             1E9, 0.1E-10, 1E-10,
             1E9, 0.4E-10, 1E-10, 
             0]

    popt, pcov = curve_fit(multi_gaussian, vel, flux, guess)
    fig,axs = plt.subplots(figsize = (4,3))
    sns.lineplot(x=vel, y=flux, ax= axs)
    sns.lineplot(x=vel, y=multi_gaussian(vel, *popt), label='Fit', ax= axs )
    sns.lineplot(x=vel, y=gaussian(vel,popt[0],popt[1],popt[2]),label='Gaussian 1', ax= axs)
    sns.lineplot(x=vel, y=gaussian(vel,popt[3],popt[4],popt[5]), label='Gaussian 2', ax= axs)
    sns.lineplot(x=vel, y=gaussian(vel,popt[6],popt[7],popt[8]),label='Gaussian 3', ax= axs)
    sns.lineplot(x=vel, y=gaussian(vel,popt[9],popt[10],popt[11]),label='Gaussian 4', ax= axs)
    axs.legend()
    plt.show()

    return popt, pcov, fig



